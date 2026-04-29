"""
Генератор презентаций на основе API ai.rt.ru
Использует LLM (Qwen) для структурирования контента и Stable Diffusion / Яндекс ART для изображений.
Входные данные: текст
Выходные форматы: .pptx, .pdf
"""

import os
import sys
import uuid
import time
import json

import textwrap
import requests
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# ──────────────────────────────────────────────
# КОНФИГУРАЦИЯ
# ──────────────────────────────────────────────

API_BASE = "https://ai.rt.ru/api/1.0"
TOKEN = "eyJhbGciOiJIUzM4NCJ9.eyJzY29wZXMiOlsibGxhbWEiLCJzZCIsInlhQXJ0Il0sInN1YiI6ImhhY2thdGhvbl8yNl8wOCIsImlhdCI6MTc3Njk0OTIzNSwiZXhwIjoxNzc3NjQwNDM1fQ.uN5ySYIA5Gxebe7TNZZQLfutNc6uwlLNPVV3kKZlyTXng_oi8_odhXa4d4VKdNpN"

# ╔══════════════════════════════════════════════════════╗
# ║  ОБНОВЛЕНО: последняя доступная версия Qwen          ║
# ╚══════════════════════════════════════════════════════╝
LLM_MODEL = "Qwen/Qwen2.5-72B-Instruct"   # ← замените на актуальную, если API поддерживает более новую

SESSION_UUID = str(uuid.uuid4())

HEADERS = {
    "Authorization": f"Bearer {TOKEN}",
    "Content-Type": "application/json",
}

# Стили оформления → цветовые схемы (R, G, B)
STYLES = {
    "корпоративный": {
        "bg":                  (15, 40, 80),
        "accent":              (0, 162, 232),
        "text":                (255, 255, 255),
        "subtext":             (180, 210, 240),
        "image_prompt_suffix": "professional corporate style, clean, blue tones",
    },
    "минималистичный": {
        "bg":                  (245, 245, 245),
        "accent":              (50, 50, 50),
        "text":                (30, 30, 30),
        "subtext":             (100, 100, 100),
        "image_prompt_suffix": "minimalist style, clean white background, simple",
    },
    "творческий": {
        "bg":                  (30, 20, 50),
        "accent":              (220, 80, 180),
        "text":                (255, 255, 255),
        "subtext":             (200, 160, 230),
        "image_prompt_suffix": "creative artistic style, vibrant colors, dynamic",
    },
    "технологичный": {
        "bg":                  (10, 20, 30),
        "accent":              (0, 220, 180),
        "text":                (200, 255, 245),
        "subtext":             (120, 200, 190),
        "image_prompt_suffix": "futuristic tech style, dark background, neon accents, digital",
    },
}

# Тоны → инструкция для LLM
TONES = {
    "нейтральный":     "Используй нейтральный, информационный тон без эмоций.",
    "формальный":      "Используй официально-деловой, строгий тон.",
    "вдохновляющий":   "Используй вдохновляющий, мотивирующий тон, апеллируй к амбициям.",
    "образовательный": "Используй понятный обучающий тон, объясняй термины, давай примеры.",
    "разговорный":     "Используй лёгкий, дружелюбный разговорный тон.",
}

# ──────────────────────────────────────────────
# ТИПЫ СЛАЙДОВ — ФИКСИРОВАННЫЕ ШАБЛОНЫ
# ──────────────────────────────────────────────
# Каждому слайду назначается один из этих типов.
# Это гарантирует, что layout не меняется случайно от запуска к запуску.
#
# ТИПЫ:
#   "title"          — титульный (большой заголовок по центру, подзаголовок)
#   "bullets"        — заголовок слева сверху + маркированный список
#   "two_col"        — два столбца: заголовок + левый текст + правый текст
#   "image_right"    — заголовок + текст слева + изображение справа
#   "highlight_block"— заголовок + буллеты слева + цветной блок-цитата справа
#   "closing"        — итоговый / спасибо (как title, но другой акцент)

SLIDE_TYPES = ["title", "bullets", "two_col", "image_right", "highlight_block", "closing"]

# Последовательность типов для N слайдов:
# Слайд 0 → всегда "title"
# Слайд N-1 → всегда "closing"
# Промежуточные → чередуются детерминированно
_MIDDLE_CYCLE = ["bullets", "highlight_block", "two_col", "image_right", "highlight_block", "bullets"]

# Ключевые слова в промпте пользователя, при которых highlight_block
# принудительно вставляется на каждый второй контентный слайд
_HIGHLIGHT_KEYWORDS = [
    "выделен", "выделить", "блок", "цитат", "важн", "акцент",
    "highlight", "quote", "callout", "featured", "featured text",
]


def _wants_highlight(user_text: str) -> bool:
    """Проверяет, просит ли пользователь выделенный блок текста."""
    low = user_text.lower()
    return any(kw in low for kw in _HIGHLIGHT_KEYWORDS)


def assign_slide_type(idx: int, total: int, force_highlight: bool = False) -> str:
    if idx == 0:
        return "title"
    if idx == total - 1:
        return "closing"
    middle_idx = idx - 1
    if force_highlight:
        # Чётные контентные слайды → highlight_block, нечётные → bullets/two_col/image_right
        if middle_idx % 2 == 0:
            return "highlight_block"
        alt_cycle = ["bullets", "two_col", "image_right"]
        return alt_cycle[(middle_idx // 2) % len(alt_cycle)]
    return _MIDDLE_CYCLE[middle_idx % len(_MIDDLE_CYCLE)]


# ──────────────────────────────────────────────
# 1. LLM — БАЗОВЫЙ ЗАПРОС
# ──────────────────────────────────────────────

def _llm_request(user_message: str, system_prompt: str, max_tokens: int = 256) -> str:
    """Базовый запрос к LLM. Возвращает строку ответа."""
    payload = {
        "uuid": SESSION_UUID,
        "chat": {
            "model": LLM_MODEL,
            "user_message": user_message,
            "contents": [
                {
                    "type": "text",
                    "text": user_message,
                    "isUrl": False,
                    "isAudioAsUrl": False,
                }
            ],
            "system_prompt": system_prompt,
            "max_new_tokens": max_tokens,
            "no_repeat_ngram_size": 15,
            "repetition_penalty": 1.1,
            "temperature": 0.2,
            "top_k": 40,
            "top_p": 0.9,
            "chat_history": [],
        },
    }
    resp = requests.post(f"{API_BASE}/llama/chat", headers=HEADERS, json=payload, timeout=60)
    if not resp.ok:
        raise RuntimeError(f"API вернул {resp.status_code}: {resp.text[:300]}")
    data = resp.json()
    if isinstance(data, list) and data:
        return data[0].get("message", {}).get("content", "").strip()
    return data.get("message", {}).get("content", "").strip()


# ──────────────────────────────────────────────
# 2. LLM — АВТООПРЕДЕЛЕНИЕ СТИЛЯ И ТОНА
# ──────────────────────────────────────────────

def detect_style_and_tone(user_text: str) -> tuple[str, str]:
    style_keys = ", ".join(STYLES.keys())
    tone_keys  = ", ".join(TONES.keys())

    system_prompt = (
        "Ты — эксперт по презентациям. "
        "Твой ответ должен содержать ТОЛЬКО валидный JSON-объект и НИЧЕГО больше. "
        "Никаких пояснений, вступлений и markdown-блоков. "
        "Первый символ ответа — '{', последний — '}'."
    )

    user_message = (
        f"Проанализируй текст запроса пользователя и выбери подходящие стиль и тон для презентации.\n\n"
        f"Доступные стили: {style_keys}\n"
        f"Доступные тоны: {tone_keys}\n\n"
        f"ВАЖНО: Если пользователь явно указывает стиль или тон "
        f"(например, 'корпоративный стиль', 'формальный тон', 'технологичный', 'разговорный' и т.п.) "
        f"— используй именно его.\n"
        f"Если явного указания нет — выбери сам, исходя из темы и контекста.\n\n"
        f"Верни JSON строго в формате:\n"
        f'{{ "style": "<один из: {style_keys}>", "tone": "<один из: {tone_keys}>" }}\n\n'
        f"ТЕКСТ ПОЛЬЗОВАТЕЛЯ:\n{user_text[:2000]}"
    )

    print("  → Определение стиля и тона через LLM...")
    raw = _llm_request(user_message, system_prompt, max_tokens=64)

    start = raw.find("{")
    end   = raw.rfind("}")
    if start != -1 and end != -1:
        raw = raw[start:end + 1]

    try:
        result = json.loads(raw)
        style  = result.get("style", "").strip().lower()
        tone   = result.get("tone",  "").strip().lower()
        if style not in STYLES:
            style = "корпоративный"
        if tone not in TONES:
            tone = "нейтральный"
        return style, tone
    except (json.JSONDecodeError, AttributeError):
        return "корпоративный", "нейтральный"


# ──────────────────────────────────────────────
# 3. LLM — СТРУКТУРИРОВАНИЕ КОНТЕНТА
# ──────────────────────────────────────────────

def generate_structure(text: str, n_slides: int, tone: str,
                       force_highlight: bool = False) -> list[dict]:
    """
    Запрашивает LLM и получает структуру презентации в JSON.
    Возвращает список слайдов: [{title, content, subtitle, left, right,
                                  highlight, highlight_label, image_prompt}, ...]
    force_highlight=True → чётные контентные слайды получают тип highlight_block.
    """
    tone_instruction = TONES.get(tone, TONES["нейтральный"])

    slide_type_list = [assign_slide_type(i, n_slides, force_highlight) for i in range(n_slides)]
    slide_hints = "\n".join(
        f"  Слайд {i+1}: тип «{t}»" for i, t in enumerate(slide_type_list)
    )

    system_prompt = (
        "Ты — эксперт по созданию презентаций. "
        "Твой ответ должен содержать ТОЛЬКО валидный JSON-массив и НИЧЕГО больше. "
        "Никаких вступлений, пояснений, комментариев и markdown-блоков. "
        "Не используй ```json, ``` или любые другие обёртки. "
        "Первый символ ответа — '[', последний — ']'."
    )

    user_message = (
        f"Создай структуру презентации ровно из {n_slides} слайдов по тексту ниже.\n"
        f"{tone_instruction}\n\n"
        f"Каждому слайду назначен тип — соблюдай его:\n{slide_hints}\n\n"
        f"Верни JSON-массив из {n_slides} объектов.\n"
        f"Структура объекта зависит от типа слайда:\n\n"
        f"  title / closing:\n"
        f'    {{"title": "...", "subtitle": "краткий подзаголовок 1 строка", "content": [], "image_prompt": "..."}}\n\n'
        f"  bullets:\n"
        f'    {{"title": "...", "content": ["Пункт 1", "Пункт 2", "Пункт 3", "Пункт 4"], "image_prompt": "..."}}\n\n'
        f"  two_col:\n"
        f'    {{"title": "...", "left": ["Пункт 1", "Пункт 2"], "right": ["Пункт 3", "Пункт 4"], "content": [], "image_prompt": "..."}}\n\n'
        f"  image_right:\n"
        f'    {{"title": "...", "content": ["Пункт 1", "Пункт 2", "Пункт 3"], "image_prompt": "..."}}\n\n'
        f"  highlight_block:\n"
        f'    {{"title": "...", "content": ["Пункт 1", "Пункт 2"], '
        f'"highlight": "Короткий яркий факт или цитата — 1–2 предложения, самое важное на этом слайде", '
        f'"highlight_label": "Знали ли вы?", "image_prompt": "..."}}\n'
        f"    (highlight — текст для визуально выделенного цветного блока; highlight_label — подпись над блоком)\n\n"
        f"Правила:\n"
        f"- Количество объектов в массиве строго равно {n_slides}\n"
        f"- image_prompt всегда на английском языке\n"
        f"- Поле «content» присутствует во всех объектах (может быть пустым [])\n"
        f"- Для слайдов типа highlight_block поля «highlight» и «highlight_label» обязательны\n\n"
        f"ТЕКСТ ДЛЯ ПРЕЗЕНТАЦИИ:\n{text[:6000]}"
    )

    print(f"  → Генерация структуры ({LLM_MODEL})...")
    raw_text = _llm_request(user_message, system_prompt, max_tokens=2048)

    if not raw_text:
        raise RuntimeError("LLM вернул пустой ответ.")

    raw_text = raw_text.strip()
    if "```" in raw_text:
        parts = raw_text.split("```")
        for part in parts:
            part = part.strip()
            if part.startswith("json"):
                part = part[4:].strip()
            if part.startswith("["):
                raw_text = part
                break
    start = raw_text.find("[")
    end   = raw_text.rfind("]")
    if start != -1 and end != -1:
        raw_text = raw_text[start:end + 1]
    raw_text = raw_text.strip()

    try:
        slides = json.loads(raw_text)
    except json.JSONDecodeError as e:
        raise RuntimeError(
            f"Не удалось разобрать JSON из ответа LLM: {e}\n"
            f"Ответ модели: {raw_text[:400]}"
        )

    if not isinstance(slides, list):
        raise RuntimeError(f"LLM вернул не массив: {type(slides)}")

    slides = slides[:n_slides]

    # Проставляем тип каждому слайду (на случай если LLM его не добавил)
    for i, s in enumerate(slides):
        s["_type"] = assign_slide_type(i, n_slides, force_highlight)

    return slides


# ──────────────────────────────────────────────
# 4. ГЕНЕРАЦИЯ ИЗОБРАЖЕНИЙ
# ──────────────────────────────────────────────

def generate_image(prompt: str, style_suffix: str, use_yandex: bool = False) -> bytes | None:
    full_prompt = f"{prompt}, {style_suffix}"
    try:
        if use_yandex:
            return _gen_yandex_art(full_prompt)
        else:
            return _gen_stable_diffusion(full_prompt)
    except Exception as e:
        print(f"    ⚠ Не удалось сгенерировать изображение: {e}")
        return None


def _gen_stable_diffusion(prompt: str) -> bytes | None:
    payload = {
        "uuid": SESSION_UUID,
        "sdImage": {
            "request": prompt,
            "seed": 42,
            "translate": True,
        },
    }
    resp = requests.post(f"{API_BASE}/sd/img", headers=HEADERS, json=payload, timeout=120)
    resp.raise_for_status()
    data = resp.json()
    img_id = data[0]["message"]["id"]
    return _download_image(img_id, "sd")


def _gen_yandex_art(prompt: str) -> bytes | None:
    payload = {
        "uuid": SESSION_UUID,
        "image": {
            "request": prompt,
            "seed": 42,
            "translate": True,
            "model": "yandex-art",
            "aspect": "16:9",
        },
    }
    resp = requests.post(f"{API_BASE}/ya/image", headers=HEADERS, json=payload, timeout=120)
    resp.raise_for_status()
    data = resp.json()
    img_id = data[0]["message"]["id"]
    return _download_image(img_id, "yaArt")


def _download_image(img_id: int, service_type: str) -> bytes:
    url = f"{API_BASE}/download"
    params = {"id": img_id, "serviceType": service_type, "imageType": "png"}
    for attempt in range(12):
        time.sleep(5)
        resp = requests.get(url, headers=HEADERS, params=params, timeout=60)
        if resp.status_code == 200 and resp.headers.get("Content-Type", "").startswith("image"):
            return resp.content
        print(f"    ⏳ Ожидание изображения ({attempt + 1}/12)...")
    raise TimeoutError("Изображение не было готово в течение 60 секунд")


# ──────────────────────────────────────────────
# 5. ПОМОЩНИКИ ДЛЯ СБОРКИ PPTX
# ──────────────────────────────────────────────

def _add_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text: str, font_size: int, bold: bool,
                 color: RGBColor, align=PP_ALIGN.LEFT,
                 word_wrap: bool = True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color


def _add_bullets(slide, left, top, width, height,
                 items: list[str], font_size: int,
                 color: RGBColor, bullet: str = "▸  ") -> None:
    if not items:
        return
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, point in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(5)
        run = para.add_run()
        run.text = f"{bullet}{point}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color


def _add_page_num(slide, idx: int, total: int, color: RGBColor):
    txb = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r   = p.add_run()
    r.text = f"{idx + 1} / {total}"
    r.font.size = Pt(11)
    r.font.color.rgb = color


def _set_shape_transparency(shape, percent: int):
    from pptx.oxml.ns import qn
    from lxml import etree
    solidFill = shape.fill._xPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is None:
        return
    alpha = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha.set("val", str(int((100 - percent) * 1000)))


# ──────────────────────────────────────────────
# 6. РЕНДЕР КАЖДОГО ТИПА СЛАЙДА
#
# Все координаты вычисляются здесь и нигде больше.
# Это гарантирует одинаковый layout при каждом запуске.
# ──────────────────────────────────────────────

# Размер слайда: 13.33 × 7.5 дюймов (широкий формат)
SW = 13.33   # ширина
SH = 7.5     # высота
M  = 0.5     # поля


def _render_title(slide, data: dict, style: dict,
                  img_bytes: bytes | None, is_closing: bool):
    """Титульный / закрывающий слайд: большой заголовок по центру."""
    txt_c = RGBColor(*style["text"])
    sub_c = RGBColor(*style["subtext"])
    acc_c = RGBColor(*style["accent"])

    title    = data.get("title", "")
    subtitle = data.get("subtitle", "") or (", ".join(data.get("content", [])) if data.get("content") else "")

    if img_bytes:
        # Фоновое изображение на всю ширину с тёмным оверлеем
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(SW), Inches(SH))
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SW), Inches(SH))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(*style["bg"])
        overlay.line.fill.background()
        _set_shape_transparency(overlay, 55)

    # Заголовок по центру
    _add_textbox(slide,
                 left=Inches(M), top=Inches(2.2),
                 width=Inches(SW - 2 * M), height=Inches(1.8),
                 text=title, font_size=44, bold=True,
                 color=txt_c, align=PP_ALIGN.CENTER)

    # Акцентная черта под заголовком — горизонтальная линия (тонкая)
    line = slide.shapes.add_shape(1, Inches(SW / 2 - 2), Inches(4.15), Inches(4), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = acc_c
    line.line.fill.background()

    if subtitle:
        _add_textbox(slide,
                     left=Inches(M), top=Inches(4.4),
                     width=Inches(SW - 2 * M), height=Inches(0.8),
                     text=subtitle, font_size=20, bold=False,
                     color=sub_c, align=PP_ALIGN.CENTER)


def _render_bullets(slide, data: dict, style: dict,
                    img_bytes: bytes | None, idx: int, total: int):
    """Стандартный слайд: заголовок вверху, маркированный список."""
    txt_c = RGBColor(*style["text"])
    sub_c = RGBColor(*style["subtext"])

    title   = data.get("title", "")
    content = data.get("content", [])

    _add_textbox(slide,
                 left=Inches(M), top=Inches(0.4),
                 width=Inches(SW - 2 * M), height=Inches(0.9),
                 text=title, font_size=30, bold=True,
                 color=txt_c)

    # Если есть изображение — текст занимает левую половину
    text_w = Inches(6.8) if img_bytes else Inches(SW - 2 * M)

    _add_bullets(slide,
                 left=Inches(M), top=Inches(1.6),
                 width=text_w, height=Inches(5.5),
                 items=content, font_size=17, color=sub_c)

    if img_bytes:
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream,
                                 Inches(7.5), Inches(1.2),
                                 Inches(5.3), Inches(5.8))

    _add_page_num(slide, idx, total, sub_c)


def _render_two_col(slide, data: dict, style: dict,
                    img_bytes: bytes | None, idx: int, total: int):
    """Двухколоночный слайд: заголовок вверху, два столбца текста."""
    txt_c = RGBColor(*style["text"])
    sub_c = RGBColor(*style["subtext"])
    acc_c = RGBColor(*style["accent"])

    title = data.get("title", "")
    left  = data.get("left",  data.get("content", [])[:3])
    right = data.get("right", data.get("content", [])[3:] or data.get("content", []))

    _add_textbox(slide,
                 left=Inches(M), top=Inches(0.4),
                 width=Inches(SW - 2 * M), height=Inches(0.9),
                 text=title, font_size=30, bold=True, color=txt_c)

    col_w = Inches(5.8)
    col_h = Inches(5.5)

    # Левый столбец
    _add_bullets(slide,
                 left=Inches(M), top=Inches(1.6),
                 width=col_w, height=col_h,
                 items=left, font_size=17, color=sub_c)

    # Разделитель
    div = slide.shapes.add_shape(1, Inches(6.66 - 0.02), Inches(1.5), Pt(2), Inches(5.7))
    div.fill.solid()
    div.fill.fore_color.rgb = acc_c
    div.line.fill.background()

    # Правый столбец
    _add_bullets(slide,
                 left=Inches(6.8), top=Inches(1.6),
                 width=col_w, height=col_h,
                 items=right, font_size=17, color=sub_c)

    _add_page_num(slide, idx, total, sub_c)


def _render_image_right(slide, data: dict, style: dict,
                         img_bytes: bytes | None, idx: int, total: int):
    """Слайд с изображением справа и буллетами слева."""
    # Идентичен bullets, но изображение обязательно
    _render_bullets(slide, data, style, img_bytes, idx, total)


def _render_highlight_block(slide, data: dict, style: dict,
                             img_bytes: bytes | None, idx: int, total: int):
    """
    Слайд с визуально выделенным блоком текста.
    Layout: заголовок вверху, буллеты слева (60%), цветной блок справа (38%).
    Блок имеет заливку цветом accent и крупный контрастный текст.
    """
    txt_c = RGBColor(*style["text"])
    sub_c = RGBColor(*style["subtext"])
    acc_c = RGBColor(*style["accent"])
    bg_c  = RGBColor(*style["bg"])

    title           = data.get("title", "")
    content         = data.get("content", [])
    highlight_text  = data.get("highlight", "")
    highlight_label = data.get("highlight_label", "Важно")

    # — Заголовок слайда
    _add_textbox(slide,
                 left=Inches(M), top=Inches(0.4),
                 width=Inches(SW - 2 * M), height=Inches(0.9),
                 text=title, font_size=30, bold=True, color=txt_c)

    # — Буллеты (левая зона, 60% ширины)
    bullet_w = Inches(7.2)
    _add_bullets(slide,
                 left=Inches(M), top=Inches(1.6),
                 width=bullet_w, height=Inches(5.5),
                 items=content, font_size=17, color=sub_c)

    # — Выделенный блок справа (38% ширины)
    block_left   = Inches(8.1)
    block_top    = Inches(1.4)
    block_width  = Inches(4.7)
    block_height = Inches(5.7)

    # Фоновый прямоугольник (цвет акцента)
    rect = slide.shapes.add_shape(
        1, block_left, block_top, block_width, block_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = acc_c
    rect.line.fill.background()

    # Скруглим углы через XML (необязательно, но красиво)
    try:
        from pptx.oxml.ns import qn
        sp_pr = rect._element.find(qn("p:spPr"))
        if sp_pr is not None:
            prstGeom = sp_pr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                prstGeom.set("prst", "roundRect")
    except Exception:
        pass

    # Подпись над блоком (маленький капс)
    label_txb = slide.shapes.add_textbox(
        block_left + Inches(0.25), block_top + Inches(0.25),
        block_width - Inches(0.5), Inches(0.4)
    )
    label_tf = label_txb.text_frame
    label_p  = label_tf.paragraphs[0]
    label_p.alignment = PP_ALIGN.LEFT
    label_run = label_p.add_run()
    label_run.text = highlight_label.upper()
    label_run.font.size  = Pt(11)
    label_run.font.bold  = True
    label_run.font.color.rgb = bg_c

    # Тонкая белая черта-разделитель под подписью
    sep = slide.shapes.add_shape(
        1,
        block_left + Inches(0.25),
        block_top + Inches(0.72),
        block_width - Inches(0.5),
        Pt(2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = bg_c
    sep.line.fill.background()

    # Сам текст выделенного блока
    text_txb = slide.shapes.add_textbox(
        block_left + Inches(0.25),
        block_top + Inches(0.9),
        block_width - Inches(0.5),
        block_height - Inches(1.1)
    )
    text_tf = text_txb.text_frame
    text_tf.word_wrap = True
    text_p  = text_tf.paragraphs[0]
    text_p.alignment = PP_ALIGN.LEFT
    text_run = text_p.add_run()
    text_run.text = highlight_text or "—"
    text_run.font.size  = Pt(19)
    text_run.font.bold  = True
    text_run.font.color.rgb = bg_c

    _add_page_num(slide, idx, total, sub_c)


# ──────────────────────────────────────────────
# 7. СБОРКА PPTX
# ──────────────────────────────────────────────

def build_pptx(slides_data: list[dict], style: dict, output_path: str,
               generate_images: bool = True, use_yandex: bool = False):
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    blank_layout = prs.slide_layouts[6]
    bg_color     = RGBColor(*style["bg"])
    total        = len(slides_data)

    for idx, slide_data in enumerate(slides_data):
        slide      = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("_type", assign_slide_type(idx, total))

        _add_bg(slide, bg_color)

        img_bytes  = None
        img_prompt = slide_data.get("image_prompt", "")
        if generate_images and img_prompt:
            print(f"  🎨 Генерация изображения для слайда {idx + 1} [{slide_type}]: «{img_prompt[:60]}»")
            img_bytes = generate_image(img_prompt, style.get("image_prompt_suffix", ""), use_yandex)

        # ── Рендер по типу ──────────────────────────────────────
        if slide_type == "title":
            _render_title(slide, slide_data, style, img_bytes, is_closing=False)

        elif slide_type == "closing":
            _render_title(slide, slide_data, style, img_bytes, is_closing=True)

        elif slide_type == "two_col":
            _render_two_col(slide, slide_data, style, img_bytes, idx, total)

        elif slide_type == "image_right":
            _render_image_right(slide, slide_data, style, img_bytes, idx, total)

        elif slide_type == "highlight_block":
            _render_highlight_block(slide, slide_data, style, img_bytes, idx, total)

        else:  # bullets (и любой неизвестный тип)
            _render_bullets(slide, slide_data, style, img_bytes, idx, total)

    prs.save(output_path)
    print(f"  ✅ PPTX сохранён: {output_path}")


# ──────────────────────────────────────────────
# 8. КОНВЕРТАЦИЯ В PDF
# ──────────────────────────────────────────────

def pptx_to_pdf(pptx_path: str, pdf_path: str, slides_data: list[dict], style: dict):
    W, H = landscape(A4)

    font_name = "Helvetica"
    for candidate in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    ]:
        if os.path.exists(candidate):
            try:
                pdfmetrics.registerFont(TTFont("CyrFont", candidate))
                font_name = "CyrFont"
            except Exception:
                pass
            break

    bold_font = font_name

    bg  = [c / 255 for c in style["bg"]]
    acc = [c / 255 for c in style["accent"]]
    txt = [c / 255 for c in style["text"]]
    sub = [c / 255 for c in style["subtext"]]

    c = rl_canvas.Canvas(pdf_path, pagesize=(W, H))

    for idx, slide_data in enumerate(slides_data):
        c.setFillColorRGB(*bg)
        c.rect(0, 0, W, H, fill=1, stroke=0)

        slide_type = slide_data.get("_type", assign_slide_type(idx, len(slides_data)))
        title      = slide_data.get("title", f"Слайд {idx + 1}")
        content    = slide_data.get("content", [])
        is_center  = slide_type in ("title", "closing")

        # Заголовок
        c.setFillColorRGB(*txt)
        c.setFont(bold_font, 30 if is_center else 22)
        title_y = H / 2 + 60 if is_center else H - 55
        x_title = W / 2 if is_center else 40
        wrapped_title = textwrap.wrap(title, width=50 if is_center else 60)
        for i, line_t in enumerate(wrapped_title):
            if is_center:
                c.drawCentredString(x_title, title_y - i * 36, line_t)
            else:
                c.drawString(x_title, title_y - i * 28, line_t)

        # Акцентная черта
        c.setStrokeColorRGB(*acc)
        c.setLineWidth(2)
        sep_y = (H / 2 + 10) if is_center else (H - 75)
        c.line(40, sep_y, W - 40, sep_y)

        # Контент
        c.setFont(font_name, 14)
        c.setFillColorRGB(*sub)

        if slide_type == "two_col":
            left_items  = slide_data.get("left",  content[:3])
            right_items = slide_data.get("right", content[3:] or content)
            col_y = H - 100
            c.line(W / 2, col_y + 10, W / 2, 40)
            for pt in left_items:
                for wl in textwrap.wrap(f"▸  {pt}", 40):
                    if col_y < 40:
                        break
                    c.drawString(40, col_y, wl)
                    col_y -= 20
                col_y -= 5
            col_y = H - 100
            for pt in right_items:
                for wl in textwrap.wrap(f"▸  {pt}", 40):
                    if col_y < 40:
                        break
                    c.drawString(W / 2 + 15, col_y, wl)
                    col_y -= 20
                col_y -= 5
        elif slide_type == "highlight_block":
            # Буллеты слева
            bullet_y = H - 100
            c.setFont(font_name, 14)
            c.setFillColorRGB(*sub)
            for pt in content:
                for wl in textwrap.wrap(f"▸  {pt}", 38):
                    if bullet_y < 40:
                        break
                    c.drawString(40, bullet_y, wl)
                    bullet_y -= 20
                bullet_y -= 5
            # Выделенный блок справа
            bx = W * 0.62
            by_top = H - 85
            bw = W * 0.36
            bh = H - 110
            c.setFillColorRGB(*acc)
            c.roundRect(bx, 30, bw, bh, 8, fill=1, stroke=0)
            # Подпись
            label = slide_data.get("highlight_label", "Важно").upper()
            c.setFillColorRGB(*bg)
            c.setFont(bold_font, 10)
            c.drawString(bx + 12, 30 + bh - 22, label)
            # Разделитель
            c.setStrokeColorRGB(*bg)
            c.setLineWidth(1)
            c.line(bx + 12, 30 + bh - 30, bx + bw - 12, 30 + bh - 30)
            # Текст блока
            hl_text = slide_data.get("highlight", "")
            c.setFont(bold_font, 16)
            hl_y = 30 + bh - 50
            for wl in textwrap.wrap(hl_text, 26):
                if hl_y < 45:
                    break
                c.drawString(bx + 12, hl_y, wl)
                hl_y -= 22
        else:
            content_y = (H / 2 - 30) if is_center else (H - 100)
            items = ([slide_data.get("subtitle", "")] if (is_center and slide_data.get("subtitle")) else []) + content
            for point in items:
                bullet = "" if is_center else "▸  "
                wrapped_p = textwrap.wrap(f"{bullet}{point}", width=75)
                for wl in wrapped_p:
                    if content_y < 40:
                        break
                    if is_center:
                        c.drawCentredString(W / 2, content_y, wl)
                    else:
                        c.drawString(50, content_y, wl)
                    content_y -= 22
                content_y -= 5

        if not is_center:
            c.setFont(font_name, 10)
            c.setFillColorRGB(*sub)
            c.drawRightString(W - 20, 15, f"{idx + 1} / {len(slides_data)}")

        c.showPage()

    c.save()
    print(f"  ✅ PDF сохранён: {pdf_path}")


# ──────────────────────────────────────────────
# 9. ИНТЕРАКТИВНЫЙ КОНСОЛЬНЫЙ МАСТЕР
# ──────────────────────────────────────────────

class C:
    RESET   = "\033[0m"
    BOLD    = "\033[1m"
    DIM     = "\033[2m"
    CYAN    = "\033[96m"
    GREEN   = "\033[92m"
    YELLOW  = "\033[93m"
    RED     = "\033[91m"
    BLUE    = "\033[94m"
    MAGENTA = "\033[95m"
    WHITE   = "\033[97m"


def clr(text: str, *codes) -> str:
    return "".join(codes) + str(text) + C.RESET


def print_header():
    print()
    print(clr("╔══════════════════════════════════════════════════════╗", C.CYAN, C.BOLD))
    print(clr("║       🎯  ГЕНЕРАТОР ПРЕЗЕНТАЦИЙ  (AI · ai.rt.ru)    ║", C.CYAN, C.BOLD))
    print(clr("╚══════════════════════════════════════════════════════╝", C.CYAN, C.BOLD))
    print()


def print_step(n: int, total: int, title: str):
    print()
    print(clr(f"── Шаг {n}/{total}: {title} ", C.BLUE, C.BOLD) + clr("─" * (40 - len(title)), C.DIM))


def choose(prompt: str, options: list[tuple[str, str]], default: int = 0) -> str:
    print(clr(prompt, C.WHITE, C.BOLD))
    for i, (key, desc) in enumerate(options, 1):
        marker       = clr(f"  [{i}]", C.CYAN, C.BOLD)
        key_str      = clr(f" {key}", C.GREEN)
        desc_str     = clr(f" — {desc}", C.DIM)
        default_mark = clr("  ← по умолчанию", C.YELLOW) if i - 1 == default else ""
        print(f"{marker}{key_str}{desc_str}{default_mark}")

    while True:
        raw = input(clr(f"\n  Введите номер [1–{len(options)}] или Enter для значения по умолчанию: ", C.YELLOW)).strip()
        if raw == "":
            return options[default][0]
        if raw.isdigit() and 1 <= int(raw) <= len(options):
            return options[int(raw) - 1][0]
        print(clr(f"  ⚠  Введите число от 1 до {len(options)}", C.RED))


def ask_int(prompt: str, min_val: int, max_val: int, default: int) -> int:
    while True:
        raw = input(clr(f"{prompt} [{min_val}–{max_val}], Enter = {default}: ", C.YELLOW)).strip()
        if raw == "":
            return default
        if raw.isdigit() and min_val <= int(raw) <= max_val:
            return int(raw)
        print(clr(f"  ⚠  Введите целое число от {min_val} до {max_val}", C.RED))


def ask_text_input() -> str:
    print(clr("\n  Введите тему или описание презентации.", C.DIM))
    print(clr("  Можно указать стиль, тон, аудиторию — нейросеть учтёт всё.", C.DIM))
    print(clr("  Завершите ввод пустой строкой.\n", C.DIM))
    lines = []
    while True:
        line = input(clr("  > ", C.CYAN))
        if line == "" and lines:
            break
        if line:
            lines.append(line)
    return "\n".join(lines)


def ask_output_name() -> str:
    raw = input(clr("\n  Имя выходного файла (без расширения), Enter = presentation: ", C.YELLOW)).strip()
    if not raw:
        return "presentation"
    safe = "".join(ch for ch in raw if ch.isalnum() or ch in "._- ")
    safe = safe.strip().replace(" ", "_")
    return safe or "presentation"


def wizard() -> dict:
    print_header()

    print_step(1, 3, "Количество слайдов")
    print(clr("  Рекомендуется: 5–10 для обычного доклада, до 20 для детального отчёта.", C.DIM))
    n_slides = ask_int("\n  Сколько слайдов сгенерировать?", 1, 20, 7)

    print_step(2, 3, "Выходной формат")
    format_options = [
        ("pptx", "PowerPoint (.pptx) — для редактирования и показа"),
        ("pdf",  "PDF (.pdf) — для печати и рассылки"),
        ("both", "Оба формата сразу"),
    ]
    out_format = choose("Выберите формат файла:", format_options)

    print_step(3, 3, "Дополнительные настройки")
    img_options = [
        ("sd",     "Stable Diffusion — быстрее, абстрактный стиль"),
        ("yandex", "Яндекс ART — качественнее, реалистичнее"),
        ("none",   "Без изображений — только текст (быстро)"),
    ]
    img_choice = choose("Генерация иллюстраций для слайдов:", img_options, default=0)
    out_name   = ask_output_name()

    print()
    print(clr("─" * 56, C.DIM))
    print(clr("  📝  Опишите презентацию", C.WHITE, C.BOLD))
    print(clr("─" * 56, C.DIM))
    source = ask_text_input()

    return {
        "source":     source,
        "n_slides":   n_slides,
        "out_format": out_format,
        "img_choice": img_choice,
        "out_name":   out_name,
    }


def print_summary(cfg: dict, style_key: str, tone_key: str, n_slides: int):
    print()
    print(clr("┌─────────────────────────────────────────────┐", C.CYAN))
    print(clr("│          📋  Параметры генерации            │", C.CYAN))
    print(clr("├─────────────────────────────────────────────┤", C.CYAN))

    slide_types_preview = " → ".join(assign_slide_type(i, n_slides) for i in range(n_slides))
    rows = [
        ("Слайдов",     str(n_slides)),
        ("Макеты",      slide_types_preview[:40]),
        ("Стиль",       style_key),
        ("Тон",         tone_key),
        ("Формат",      cfg["out_format"].upper()),
        ("Изображения", {"sd": "Stable Diffusion", "yandex": "Яндекс ART", "none": "Нет"}[cfg["img_choice"]]),
        ("Файл",        cfg["out_name"] + ".[pptx/pdf]"),
    ]
    for label, value in rows:
        print(clr(f"│  {label:<14}", C.CYAN) + clr(f"{value:<30}", C.GREEN) + clr("│", C.CYAN))
    print(clr("└─────────────────────────────────────────────┘", C.CYAN))
    print()


def confirm(prompt: str = "Начать генерацию? [Y/n]: ") -> bool:
    ans = input(clr(prompt, C.YELLOW, C.BOLD)).strip().lower()
    return ans in ("", "y", "да", "д", "yes")


# ──────────────────────────────────────────────
# ТОЧКА ВХОДА
# ──────────────────────────────────────────────

def main():
    try:
        cfg = wizard()
    except KeyboardInterrupt:
        print(clr("\n\n  Отменено пользователем.", C.YELLOW))
        sys.exit(0)

    print(clr("\n🔍 Анализ запроса...", C.BLUE, C.BOLD))
    try:
        style_key, tone_key = detect_style_and_tone(cfg["source"])
    except Exception as e:
        print(clr(f"  ⚠ Не удалось определить стиль/тон: {e}. Используются значения по умолчанию.", C.YELLOW))
        style_key, tone_key = "корпоративный", "нейтральный"

    force_highlight = _wants_highlight(cfg["source"])
    if force_highlight:
        print(clr("   ✦ Обнаружен запрос на выделенный блок — включён режим highlight_block", C.MAGENTA))

    print(clr(f"   Стиль: {style_key}  |  Тон: {tone_key}", C.DIM))

    print_summary(cfg, style_key, tone_key, cfg["n_slides"])

    if not confirm():
        print(clr("  Отменено.", C.DIM))
        sys.exit(0)

    style      = STYLES[style_key]
    gen_images = cfg["img_choice"] != "none"
    use_yandex = cfg["img_choice"] == "yandex"
    pptx_path  = f"{cfg['out_name']}.pptx"
    pdf_path   = f"{cfg['out_name']}.pdf"

    # ── Генерация структуры ──────────────────────────────────────
    print(clr("\n🤖 [1/2] Генерация структуры через LLM...", C.BLUE, C.BOLD))
    try:
        slides_data = generate_structure(cfg["source"], cfg["n_slides"], tone_key,
                                         force_highlight=force_highlight)
    except Exception as e:
        print(clr(f"\n❌ Ошибка LLM: {e}", C.RED))
        sys.exit(1)

    print(clr(f"\n   Структура готова — {len(slides_data)} слайдов:", C.GREEN))
    for i, s in enumerate(slides_data):
        t    = s.get("_type", "?")
        icon = "🏷 " if i == 0 else ("🏁" if i == len(slides_data) - 1 else f"{i:>2}.")
        print(clr(f"   {icon} [{t:<12}] {s.get('title', '—')}", C.DIM))

    # ── Сборка файлов ────────────────────────────────────────────
    print(clr("\n🎨 [2/2] Сборка презентации...", C.BLUE, C.BOLD))

    if cfg["out_format"] in ("pptx", "both"):
        build_pptx(slides_data, style, pptx_path,
                   generate_images=gen_images, use_yandex=use_yandex)

    if cfg["out_format"] in ("pdf", "both"):
        pptx_to_pdf(None, pdf_path, slides_data, style)

    # ── Итог ─────────────────────────────────────────────────────
    print()
    print(clr("╔══════════════════════════════════════════════════╗", C.GREEN, C.BOLD))
    print(clr("║              ✅  Готово!                         ║", C.GREEN, C.BOLD))
    print(clr("╠══════════════════════════════════════════════════╣", C.GREEN, C.BOLD))
    if cfg["out_format"] in ("pptx", "both"):
        abs_path = os.path.abspath(pptx_path)
        print(clr(f"║  📊 {abs_path:<44}║", C.GREEN))
    if cfg["out_format"] in ("pdf", "both"):
        abs_path = os.path.abspath(pdf_path)
        print(clr(f"║  📄 {abs_path:<44}║", C.GREEN))
    print(clr("╚══════════════════════════════════════════════════╝", C.GREEN, C.BOLD))
    print()


if __name__ == "__main__":
    main()
