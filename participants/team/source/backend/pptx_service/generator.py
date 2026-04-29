# backend/pptx_service/generator.py
import io
import requests
from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

from .style_catalog import StyleCatalog

CACHE_DIR = Path("./tmp/cache")

class PPTXGenerator:
    def __init__(self, style_preset: str = "corporate", styles_catalog_folder: str = "styles"):
        self.catalog = StyleCatalog(styles_catalog_folder)
        self.style_data = self.catalog.get_style_data(style_preset) or self.catalog.presets["corporate"]
        self.style_name = style_preset
    
    def get_style_info(self) -> Dict[str, Any]:
        return self.style_data
    
    def generate(self, data: Dict[str, Any]) -> io.BytesIO:
        prs = Presentation()
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = data.get("title", "Презентация")
        subtitle.text = f"Сгенерировано AI • {len(data.get('slides', []))} слайдов"
        
        for slide_data in data.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_data.get("title", "Без названия")
            
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            bullets = slide_data.get("bullets", slide_data.get("points", []))
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(self.style_data.get("fonts", {}).get("body_size_pt", 18))
                p.font.name = self.style_data.get("fonts", {}).get("font_family", "Arial")
            
            image_url = slide_data.get("image_url")
            if image_url:
                self._add_image_to_slide(slide, image_url)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output
    
    def _add_image_to_slide(self, slide, image_url: str):
        """Синхронная версия — без asyncio и nest_asyncio"""
        try:
            img_stream = None
            
            if image_url.startswith("/"):
                file_path = CACHE_DIR / image_url.replace("/images/", "")
                if file_path.exists():
                    img_stream = io.BytesIO(file_path.read_bytes())
            else:
                resp = requests.get(image_url, timeout=10)
                if resp.status_code == 200:
                    img_stream = io.BytesIO(resp.content)
            
            if img_stream:
                img = Image.open(img_stream)
                img.thumbnail((500, 500), Image.Resampling.LANCZOS)
                
                output = io.BytesIO()
                img.save(output, format="PNG")
                output.seek(0)
                
                slide.shapes.add_picture(output, Inches(5.5), Inches(2), width=Inches(4))
                
        except Exception as e:
            print(f"⚠️ Не удалось добавить изображение: {e}")
