import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Маппинг стандартных layout'ов PowerPoint
LAYOUT_MAP = {
    "title_slide": 0,    # Титульный
    "title_content": 1,  # Заголовок + текст
    "section_header": 2, # Разделитель
    "two_columns": 3,    # Сравнение/Две колонки
    "title_only": 5,     # Только заголовок
    "blank": 6           # Пустой
}

def build_pptx(slides_data: dict, output_path: str, image_dir: str = None) -> str:
    """
    Создаёт презентацию PPTX с улучшенным стилем.
    """
    print(f"📊 Сборка презентации: {output_path}")
    
    # Создаем директорию, если её нет
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation()
    
    # Широкоформатный формат 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slides_list = slides_data.get("slides", [])
    
    if not slides_list:
        _add_empty_warning(prs)
    
    for i, info in enumerate(slides_list):
        title_text = info.get("title", f"Слайд {i+1}")
        print(f"  📄 [{i+1}/{len(slides_list)}] {title_text[:40]}...")

        layout_key = info.get("layout", "title_content")
        layout_idx = LAYOUT_MAP.get(layout_key, 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        image_name = info.get("image_path")
        full_img_path = os.path.join(image_dir, image_name) if image_dir and image_name else None
        has_image = full_img_path and os.path.exists(full_img_path)

        # 1. Заголовок (фиксируем положение)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            
            # Устанавливаем координаты и растягиваем на почти весь слайд
            slide.shapes.title.left = Inches(0.5)
            slide.shapes.title.top = Inches(0.3)
            slide.shapes.title.width = Inches(12.0) # Растягиваем вширь!
            slide.shapes.title.height = Inches(1.0) # Даем высоту для одной-двух строк
            
            _style_title(slide.shapes.title)

        # 2. Основной текст
        content = info.get("content", "")
        if content:
            body_shape = _get_body_placeholder(slide)
            if body_shape:
                # ЗАДАЕМ КООРДИНАТЫ ТЕКСТА
                body_shape.top = Inches(1.5) # Начинаем строго под заголовком
                body_shape.left = Inches(0.5) 
                
                if has_image:
                    body_shape.width = Inches(6.2) # Ограничиваем ширину, чтобы не лез под картинку
                else:
                    body_shape.width = Inches(12.0) # На весь слайд, если нет фото
                
                _fill_text_frame(body_shape.text_frame, content)

        # 3. Картинка (справа)
        if has_image:
            _add_image_to_slide(slide, full_img_path)
        elif image_name:
            print(f"  ⚠️ Изображение не найдено: {full_img_path}")

    prs.save(output_path)
    print(f"✅ Готово! Размер: {os.path.getsize(output_path)} байт")
    return output_path

# --- Вспомогательные функции для "чистого" кода ---

def _style_title(title_shape):
    """Применяет фирменный стиль к заголовку."""
    tf = title_shape.text_frame
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(10, 50, 150) # Темно-синий

def _fill_text_frame(tf, text):
    """Заполняет текстовый блок с поддержкой списков."""
    tf.clear()
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clean_line = line.strip()
        
        # Если строка — элемент списка
        if clean_line.startswith(('-', '•', '*')):
            p.level = 1
            p.text = clean_line[1:].strip()
        else:
            p.text = clean_line
            
        p.font.size = Pt(20)
        p.space_after = Pt(10)

def _add_image_to_slide(slide, img_path):
    """Аккуратно вставляет картинку в правую часть слайда."""
    # Стандартные координаты для правой части 16:9
    left = Inches(8.5)
    top = Inches(1.5)
    width = Inches(4.2) 
    slide.shapes.add_picture(img_path, left, top, width=width)

def _get_body_placeholder(slide):
    """Находит placeholder для текста (пропуская заголовок)."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx > 0:
            return shape
    return None
def _add_empty_warning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Данные отсутствуют"