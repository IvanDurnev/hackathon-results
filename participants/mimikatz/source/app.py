import io
import json
import os
import random
import re
import tempfile
import uuid
from http.client import IncompleteRead
from pathlib import Path
from time import sleep, time
from typing import Any, Dict, List, Optional

import fitz  # PyMuPDF
import requests
from docx import Document
from dotenv import load_dotenv
from flask import Flask, jsonify, render_template, request, send_file
# Для работы с pydub нужен ffmpeg: brew install ffmpeg или apt install ffmpeg
from pydub import AudioSegment
import speech_recognition as sr
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
from werkzeug.utils import secure_filename

load_dotenv(override=True)


app = Flask(__name__)

API_TOKEN = os.getenv("HACKATHON_API_TOKEN", "")
AI_BASE_URL = os.getenv("AI_BASE_URL", "https://ai.rt.ru/api/1.0")
LLM_URL = os.getenv("LLM_URL", f"{AI_BASE_URL}/llama/chat")
IMAGE_GEN_URL = os.getenv("IMAGE_GEN_URL", f"{AI_BASE_URL}/ya/image")
IMAGE_DOWNLOAD_URL = os.getenv("IMAGE_DOWNLOAD_URL", f"{AI_BASE_URL}/download")
REQUEST_TIMEOUT = int(os.getenv("REQUEST_TIMEOUT", "90"))
IMAGE_REQUEST_TIMEOUT = int(os.getenv("IMAGE_REQUEST_TIMEOUT", "180"))
IMAGE_DOWNLOAD_TIMEOUT = int(os.getenv("IMAGE_DOWNLOAD_TIMEOUT", "120"))
IMAGE_GEN_RETRIES = int(os.getenv("IMAGE_GEN_RETRIES", "2"))
IMAGE_GEN_BACKOFF = float(os.getenv("IMAGE_GEN_BACKOFF", "2.0"))
IMAGE_PREVIEW_RETRIES = int(os.getenv("IMAGE_PREVIEW_RETRIES", "5"))
IMAGE_PREVIEW_RETRY_DELAY = float(os.getenv("IMAGE_PREVIEW_RETRY_DELAY", "1.2"))
LLM_MODEL = os.getenv("LLM_MODEL", "Qwen/Qwen2.5-72B-Instruct")

ALLOWED_EXTENSIONS = {"pdf", "docx", "txt"}
MAX_CONTEXT_CHARS = int(os.getenv("MAX_CONTEXT_CHARS", "12000"))

STYLE_PRESETS: Dict[str, Dict[str, Any]] = {
    "professional": {
        "bg": RGBColor(245, 247, 250),
        "title": RGBColor(33, 37, 41),
        "body": RGBColor(55, 65, 81),
        "title_size": 32,
        "body_size": 18,
        "font_family": "Calibri",
        "layout": "standard",
    },
    "dark": {
        "bg": RGBColor(22, 27, 34),
        "title": RGBColor(230, 237, 243),
        "body": RGBColor(201, 209, 217),
        "title_size": 32,
        "body_size": 18,
        "font_family": "Calibri",
        "layout": "standard",
    },
    "minimal": {
        "bg": RGBColor(255, 255, 255),
        "title": RGBColor(17, 24, 39),
        "body": RGBColor(75, 85, 99),
        "title_size": 32,
        "body_size": 16,
        "font_family": "Calibri",
        "layout": "standard",
    },
}

BACKGROUND_OPTIONS = {
    "white": RGBColor(255, 255, 255),
    "light_gray": RGBColor(240, 240, 240),
    "dark_blue": RGBColor(20, 30, 50),
    "gradient_blue": RGBColor(100, 150, 200),  # placeholder, gradients not supported easily
    "pastel_pink": RGBColor(255, 200, 220),
    "green": RGBColor(200, 255, 200),
    "black": RGBColor(0, 0, 0),
}

FONT_OPTIONS = {
    "Arial": "Arial",
    "Calibri": "Calibri",
    "Times New Roman": "Times New Roman",
    "Georgia": "Georgia",
    "Verdana": "Verdana",
    "Tahoma": "Tahoma",
    "Comic Sans MS": "Comic Sans MS",
    "Impact": "Impact",
    "Helvetica": "Helvetica",
    "Courier New": "Courier New",
}

LAYOUT_OPTIONS = {
    "centered": "centered",
    "left": "left",
    "right": "right",
    "asymmetric": "asymmetric",
}

IMAGE_POSITION_OPTIONS = {
    "center": "center",
    "left": "left",
    "right": "right",
    "top": "top",
    "bottom": "bottom",
}

# расширенные стили Slidesgo с текстурами и деталями
SLIDESGO_STYLES_DESIGN: Dict[str, Dict[str, Any]] = {
    "lifestyle": {
        "bg": "white",
        "title": RGBColor(45, 41, 38),
        "body": RGBColor(100, 95, 89),
        "accent": RGBColor(214, 108, 62),
        "title_size": 40,
        "body_size": 18,
        "font_family": "Helvetica",
        "layout": "asymmetric",
        "image_position": "right",
    },
    "editorialportrait": {
        "bg": RGBColor(255, 255, 255),
        "title": RGBColor(30, 30, 30),
        "body": RGBColor(80, 80, 80),
        "accent": RGBColor(200, 200, 200),
        "title_size": 36,
        "body_size": 16,
        "font_family": "Arial",
        "layout": "centered",
        "image_position": "center",
    },
    "illustration": {
        "bg": RGBColor(248, 245, 240),
        "title": RGBColor(50, 45, 40),
        "body": RGBColor(100, 95, 90),
        "accent": RGBColor(255, 180, 90),
        "title_size": 38,
        "body_size": 17,
        "font_family": "Courier New",
        "layout": "creative",
        "image_position": "left",
    },
    "3d": {
        "bg": RGBColor(20, 30, 50),
        "title": RGBColor(255, 255, 255),
        "body": RGBColor(200, 210, 220),
        "accent": RGBColor(100, 200, 255),
        "title_size": 42,
        "body_size": 18,
        "font_family": "Arial",
        "layout": "bold",
        "image_position": "right",
    },
    "graphite": {
        "bg": RGBColor(240, 240, 240),
        "title": RGBColor(50, 50, 50),
        "body": RGBColor(100, 100, 100),
        "accent": RGBColor(80, 80, 80),
        "title_size": 36,
        "body_size": 16,
        "font_family": "Courier New",
        "layout": "minimalist",
        "image_position": "center",
    },
    "watercolor": {
        "bg": RGBColor(245, 250, 248),
        "title": RGBColor(60, 100, 120),
        "body": RGBColor(100, 130, 150),
        "accent": RGBColor(150, 100, 200),
        "title_size": 40,
        "body_size": 17,
        "font_family": "Georgia",
        "layout": "artistic",
        "image_position": "left",
    },
    "3dcartoon": {
        "bg": RGBColor(255, 245, 230),
        "title": RGBColor(100, 60, 40),
        "body": RGBColor(120, 90, 70),
        "accent": RGBColor(255, 150, 70),
        "title_size": 42,
        "body_size": 18,
        "font_family": "Comic Sans MS",
        "layout": "playful",
        "image_position": "right",
    },
    "anime": {
        "bg": RGBColor(255, 240, 245),
        "title": RGBColor(80, 40, 120),
        "body": RGBColor(120, 80, 140),
        "accent": RGBColor(255, 100, 150),
        "title_size": 40,
        "body_size": 18,
        "font_family": "Impact",
        "layout": "dynamic",
        "image_position": "center",
    },
    "moderncollage": {
        "bg": RGBColor(220, 225, 230),
        "title": RGBColor(40, 50, 70),
        "body": RGBColor(80, 90, 110),
        "accent": RGBColor(100, 150, 200),
        "title_size": 38,
        "body_size": 16,
        "font_family": "Verdana",
        "layout": "magazine",
        "image_position": "mixed",
    },
    "futuristictech": {
        "bg": RGBColor(15, 30, 50),
        "title": RGBColor(0, 200, 255),
        "body": RGBColor(100, 220, 255),
        "accent": RGBColor(0, 255, 200),
        "title_size": 44,
        "body_size": 18,
        "font_family": "Impact",
        "layout": "tech",
        "image_position": "right",
    },
    "blackandwhite": {
        "bg": RGBColor(255, 255, 255),
        "title": RGBColor(0, 0, 0),
        "body": RGBColor(50, 50, 50),
        "accent": RGBColor(100, 100, 100),
        "title_size": 36,
        "body_size": 16,
        "font_family": "Times New Roman",
        "layout": "classic",
        "image_position": "center",
    },
    "inspirational": {
        "bg": RGBColor(250, 250, 250),
        "title": RGBColor(100, 100, 150),
        "body": RGBColor(120, 120, 170),
        "accent": RGBColor(200, 150, 100),
        "title_size": 42,
        "body_size": 18,
        "font_family": "Georgia",
        "layout": "inspiring",
        "image_position": "left",
    },
    "neonminimal": {
        "bg": RGBColor(10, 10, 10),
        "title": RGBColor(255, 0, 255),
        "body": RGBColor(200, 200, 200),
        "accent": RGBColor(0, 255, 255),
        "title_size": 40,
        "body_size": 16,
        "font_family": "Arial",
        "layout": "minimal",
        "image_position": "right",
    },
    "futurewave": {
        "bg": RGBColor(0, 20, 40),
        "title": RGBColor(255, 100, 200),
        "body": RGBColor(150, 200, 255),
        "accent": RGBColor(100, 255, 100),
        "title_size": 38,
        "body_size": 17,
        "font_family": "Helvetica",
        "layout": "wave",
        "image_position": "center",
    },
    "pixelart3d": {
        "bg": RGBColor(200, 200, 255),
        "title": RGBColor(50, 50, 100),
        "body": RGBColor(100, 100, 150),
        "accent": RGBColor(255, 100, 100),
        "title_size": 36,
        "body_size": 16,
        "font_family": "Courier New",
        "layout": "pixel",
        "image_position": "left",
    },
    "3dplastiline": {
        "bg": RGBColor(255, 255, 200),
        "title": RGBColor(100, 50, 0),
        "body": RGBColor(150, 100, 50),
        "accent": RGBColor(255, 200, 0),
        "title_size": 42,
        "body_size": 18,
        "font_family": "Comic Sans MS",
        "layout": "plastic",
        "image_position": "right",
    },
    "highdefinition": {
        "bg": RGBColor(240, 248, 255),
        "title": RGBColor(0, 100, 200),
        "body": RGBColor(50, 150, 250),
        "accent": RGBColor(0, 200, 100),
        "title_size": 40,
        "body_size": 18,
        "font_family": "Calibri",
        "layout": "hd",
        "image_position": "center",
    },
    "historical": {
        "bg": RGBColor(245, 245, 220),
        "title": RGBColor(100, 50, 0),
        "body": RGBColor(120, 80, 40),
        "accent": RGBColor(200, 150, 50),
        "title_size": 38,
        "body_size": 17,
        "font_family": "Times New Roman",
        "layout": "vintage",
        "image_position": "left",
    },
    "darkconceptual": {
        "bg": RGBColor(20, 20, 20),
        "title": RGBColor(255, 255, 255),
        "body": RGBColor(180, 180, 180),
        "accent": RGBColor(255, 0, 0),
        "title_size": 44,
        "body_size": 18,
        "font_family": "Arial",
        "layout": "dark",
        "image_position": "right",
    },
    "surreal": {
        "bg": RGBColor(50, 0, 50),
        "title": RGBColor(255, 200, 255),
        "body": RGBColor(200, 150, 200),
        "accent": RGBColor(255, 100, 255),
        "title_size": 40,
        "body_size": 18,
        "font_family": "Fantasy",
        "layout": "surreal",
        "image_position": "center",
    },
}

WORK_DIR = Path(tempfile.gettempdir()) / "rostel_pptx"
WORK_DIR.mkdir(parents=True, exist_ok=True)
GENERATED_IMAGE_DIR = Path(app.root_path) / "static" / "generated_images"
GENERATED_IMAGE_DIR.mkdir(parents=True, exist_ok=True)
API_SESSION = requests.Session()
API_SESSION.trust_env = False
API_RETRY_TOTAL = int(os.getenv("API_RETRY_TOTAL", "2"))
API_RETRY_BACKOFF = float(os.getenv("API_RETRY_BACKOFF", "1.2"))
retry_strategy = Retry(
    total=API_RETRY_TOTAL,
    connect=API_RETRY_TOTAL,
    read=API_RETRY_TOTAL,
    status=API_RETRY_TOTAL,
    backoff_factor=API_RETRY_BACKOFF,
    status_forcelist=[429, 500, 502, 503, 504],
    allowed_methods=frozenset({"GET", "POST"}),
)
adapter = HTTPAdapter(max_retries=retry_strategy)
API_SESSION.mount("https://", adapter)
API_SESSION.mount("http://", adapter)
PUBLIC_MINUTES_PER_SLIDE = float(os.getenv("PUBLIC_MINUTES_PER_SLIDE", "0.5"))
SELF_STUDY_MINUTES_PER_SLIDE = float(os.getenv("SELF_STUDY_MINUTES_PER_SLIDE", "1.0"))
SLIDESGO_STYLE_CATEGORY_TITLES = {
    "all": "Все стили",
    "photo": "Фотографии",
    "illustration": "Иллюстрации",
    "3d": "3D",
}
SLIDESGO_STYLE_TITLE_OVERRIDES = {
    "lifestyle": "Лайфстайл",
    "editorialportrait": "Редакционный портрет",
    "illustration": "Иллюстрация",
    "3d": "3D",
    "graphite": "Графит",
    "watercolor": "Акварель",
    "3dcartoon": "3D-мультфильм",
    "anime": "Аниме",
    "moderncollage": "Современный коллаж",
    "futuristictech": "Футуристичные технологии",
    "blackandwhite": "Черно-белый",
    "inspirational": "Вдохновляющий",
    "neonminimal": "Неоновый минимализм",
    "futurewave": "Волна будущего",
    "pixelart3d": "Пиксель-арт 3D",
    "3dplastiline": "3D-пластилин",
    "highdefinition": "Высокая четкость",
    "historical": "Исторический",
    "darkconceptual": "Темный концептуальный",
    "surreal": "Сюрреализм",
}
SLIDESGO_STYLE_CATEGORY_MAP = {
    "lifestyle": "photo",
    "editorialportrait": "photo",
    "illustration": "illustration",
    "3d": "3d",
    "graphite": "illustration",
    "watercolor": "illustration",
    "3dcartoon": "3d",
    "anime": "illustration",
    "moderncollage": "illustration",
    "futuristictech": "illustration",
    "blackandwhite": "photo",
    "inspirational": "photo",
    "neonminimal": "illustration",
    "futurewave": "illustration",
    "pixelart3d": "illustration",
    "3dplastiline": "3d",
    "highdefinition": "illustration",
    "historical": "photo",
    "darkconceptual": "illustration",
    "surreal": "illustration",
}
SLIDESGO_STYLE_ITEMS = [
    {
        "key": style_key,
        "title": SLIDESGO_STYLE_TITLE_OVERRIDES.get(style_key, style_key),
        "description": "",
        "category": SLIDESGO_STYLE_CATEGORY_MAP.get(style_key, "all"),
        "config": style_config.copy(),
    }
    for style_key, style_config in SLIDESGO_STYLES_DESIGN.items()
]


def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_chunks: List[str] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text_chunks.append(page.get_text("text"))
    return "\n".join(text_chunks)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_text(uploaded_file) -> str:
    filename = secure_filename(uploaded_file.filename or "")
    extension = filename.rsplit(".", 1)[1].lower() if "." in filename else ""
    content = uploaded_file.read()

    if extension == "pdf":
        return extract_text_from_pdf(content)
    if extension == "docx":
        return extract_text_from_docx(content)
    if extension == "txt":
        return content.decode("utf-8", errors="ignore")
    return ""


def llm_headers() -> Dict[str, str]:
    if not API_TOKEN:
        raise RuntimeError("Missing HACKATHON_API_TOKEN in environment.")
    return {"Authorization": f"Bearer {API_TOKEN}", "Content-Type": "application/json"}


def presentation_mode_label(presentation_type: str) -> str:
    return "публичная" if presentation_type == "public" else "для самостоятельного ознакомления"


def presentation_mode_guidance(presentation_type: str) -> str:
    if presentation_type == "public":
        return (
            "Формат презентации: публичное выступление. "
            "На слайдах должно быть мало текста, короткие тезисы, больше визуальных образов. "
            "Для каждого слайда старайся давать выразительный image_prompt."
        )
    return (
        "Формат презентации: самостоятельное ознакомление. "
        "На слайдах должно быть больше полезного текста и пояснений, изображения нужны реже и только когда усиливают понимание."
    )


def calculate_slide_count(duration_minutes: int, presentation_type: str) -> int:
    minutes_per_slide = PUBLIC_MINUTES_PER_SLIDE if presentation_type == "public" else SELF_STUDY_MINUTES_PER_SLIDE
    estimated_slides = max(1, round(duration_minutes / minutes_per_slide))
    return max(3, min(estimated_slides, 20))


def normalize_structure(raw_data: Dict[str, Any], fallback_count: int) -> Dict[str, List[Dict[str, str]]]:
    slides: List[Dict[str, str]] = []
    raw_slides = raw_data.get("slides", [])
    if not isinstance(raw_slides, list):
        raw_slides = []

    for index, slide in enumerate(raw_slides[:fallback_count], start=1):
        if not isinstance(slide, dict):
            continue
        title = str(slide.get("title", f"Слайд {index}")).strip() or f"Слайд {index}"
        raw_content = str(slide.get("content", "")).strip()
        raw_data_blocks = slide.get("data_blocks", [])
        if isinstance(raw_data_blocks, list) and raw_data_blocks:
            data_blocks = raw_data_blocks
            content = raw_content
        else:
            data_blocks = extract_data_blocks(raw_content)
            content = strip_data_blocks(raw_content).strip()
        if not content and not data_blocks:
            content = "- Добавьте содержание"
        image_prompt = str(slide.get("image_prompt", "")).strip()
        image_id = str(slide.get("image_id", "")).strip()
        image_url = str(slide.get("image_url", "")).strip()
        image_enabled = bool(slide.get("image_enabled", True)) if image_prompt else False
        try:
            image_count = max(1, min(int(slide.get("image_count", 1)), 3))
        except (TypeError, ValueError):
            image_count = 1
        slides.append(
            {
                "title": title,
                "content": content,
                "image_prompt": image_prompt,
                "image_id": image_id,
                "image_url": image_url,
                "data_blocks": data_blocks,
                "image_enabled": image_enabled,
                "image_count": image_count,
            }
        )

    while len(slides) < fallback_count:
        i = len(slides) + 1
        slides.append(
            {
                "title": f"Слайд {i}",
                "content": "- Добавьте содержание",
                "image_prompt": "Абстрактная иллюстрация по теме",
                "image_id": "",
                "image_url": "",
                "data_blocks": [],
                "image_enabled": True,
                "image_count": 1,
            }
        )

    return {"slides": slides}


def get_slidesgo_styles() -> Dict[str, Any]:
    styles_by_category: Dict[str, List[Dict[str, Any]]] = {"all": []}
    for style in SLIDESGO_STYLE_ITEMS:
        category = style.get("category", "all")
        styles_by_category.setdefault(category, []).append(style)
        styles_by_category["all"].append(style)

    categories = [
        {"key": "all", "title": SLIDESGO_STYLE_CATEGORY_TITLES.get("all", "Все стили")}
    ]
    for category in ["photo", "illustration", "3d"]:
        if category in styles_by_category:
            categories.append({"key": category, "title": SLIDESGO_STYLE_CATEGORY_TITLES.get(category, category)})

    return {"categories": categories, "styles": styles_by_category}


def get_slidesgo_style_by_key(key: str) -> Optional[Dict[str, Any]]:
    if not key:
        return None
    for style in SLIDESGO_STYLE_ITEMS:
        if style.get("key") == key:
            return style
    return None


def humanize_template_name(stem: str) -> str:
    parts = [part for part in stem.replace("-", "_").split("_") if part and part != "style"]
    if not parts:
        return "Шаблон"
    return " ".join(parts).capitalize()


@app.route("/api/pptx_templates")
def api_pptx_templates():
    templates_dir = Path(app.root_path) / "static" / "templates"
    previews_dir = templates_dir / "previews"
    if not templates_dir.exists():
        return jsonify({"templates": []})

    templates: List[Dict[str, str]] = []
    for template_path in sorted(templates_dir.glob("style_*.pptx")):
        stem = template_path.stem
        style_key = stem.removeprefix("style_")
        preview_url = ""
        for ext in ("png", "jpg", "jpeg", "webp"):
            preview_path = previews_dir / f"{stem}.{ext}"
            if preview_path.exists():
                preview_url = f"/static/templates/previews/{preview_path.name}"
                break
        templates.append(
            {
                "key": style_key,
                "title": humanize_template_name(stem),
                "file_name": template_path.name,
                "preview_url": preview_url,
            }
        )
    return jsonify({"templates": templates})


@app.route("/api/voice-input", methods=["POST"])
def api_voice_input():
    if "file" not in request.files:
        return jsonify({"error": "Аудиофайл не найден."}), 400

    audio_file = request.files["file"]
    if not audio_file.filename:
        return jsonify({"error": "Пустое имя файла."}), 400

    # Determine format from filename
    file_ext = audio_file.filename.rsplit('.', 1)[-1].lower() if '.' in audio_file.filename else 'webm'
    if file_ext not in ['webm', 'mp4', 'ogg', 'wav', 'm4a', 'flac']:
        file_ext = 'webm'  # default

    # Save to temp file
    with tempfile.NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_file:
        temp_audio_path = temp_file.name
        audio_file.save(temp_audio_path)

    try:
        # Use pydub to convert any audio format to WAV
        audio_segment = AudioSegment.from_file(temp_audio_path)
        wav_io = io.BytesIO()
        audio_segment.export(wav_io, format="wav")
        wav_io.seek(0)

        recognizer = sr.Recognizer()
        with sr.AudioFile(wav_io) as source:
            audio_data = recognizer.record(source)
        text = recognizer.recognize_google(audio_data, language="ru-RU")
        return jsonify({"text": text})
    except sr.UnknownValueError:
        return jsonify({"error": "Речь не распознана, попробуйте еще раз."}), 400
    except sr.RequestError as err:
        return jsonify({"error": f"Ошибка сервиса распознавания: {err}"}), 502
    except Exception as err:
        return jsonify({"error": f"Ошибка обработки аудио: {err}"}), 500
    finally:
        # Clean up temp file
        if os.path.exists(temp_audio_path):
            os.unlink(temp_audio_path)


def build_system_prompt(slide_count: int, tone: str, style: str, presentation_type: str, slidesgo_style: str = "") -> str:
    prompt = (
        "Ты — профессиональный аналитик и дизайнер презентаций. "
        f"Сформируй структуру презентации из {slide_count} слайдов. "
        f"Тон: {tone}. Визуальный стиль: {style}. "
    )
    if slidesgo_style:
        prompt += f"Используй шаблон Slidesgo: {slidesgo_style}. "
    prompt += (
        f"{presentation_mode_guidance(presentation_type)} "
        "Верни строго JSON формата: "
        '{"slides":[{"title":"...","content":"- тезис 1\\n- тезис 2\\n- тезис 3","image_prompt":"..."}]}.'
        "В тексте слайда можно вставлять визуальные блоки данных внутри тегов [DATA]...[/DATA]. "
        "Поддерживаемые типы: bar, line, pie, table. "
        "Пример: [DATA] {\"type\": \"bar\", \"title\": \"Продажи\", \"labels\": [\"Q1\", \"Q2\"], \"values\": [400,500]} [/DATA]. "
        "Для таблицы: [DATA] {\"type\": \"table\", \"rows\": [[\"Заголовок1\", \"Заголовок2\"], [\"Данные1\", \"Данные2\"]]} [/DATA]. "
        "Если JSON внутри [DATA] поврежден, игнорируй визуальный элемент и возвращай только текст.")
    return prompt


def extract_json_object(text: str) -> Dict[str, Any]:
    if not text:
        raise RuntimeError("LLM вернула пустой ответ.")

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    fenced_match = re.search(r"```json\s*(\{.*?\})\s*```", text, flags=re.DOTALL | re.IGNORECASE)
    if fenced_match:
        return json.loads(fenced_match.group(1))

    object_match = re.search(r"(\{.*\})", text, flags=re.DOTALL)
    if object_match:
        return json.loads(object_match.group(1))

    raise RuntimeError("LLM вернула ответ, который не удалось распарсить как JSON.")


def strip_data_blocks(text: str) -> str:
    return re.sub(r"\[DATA\].*?\[/DATA\]", "", text, flags=re.DOTALL | re.IGNORECASE).strip()


def extract_data_blocks(text: str) -> List[Dict[str, Any]]:
    blocks: List[Dict[str, Any]] = []
    for raw_match in re.findall(r"\[DATA\](.*?)\[/DATA\]", text, flags=re.DOTALL | re.IGNORECASE):
        payload = raw_match.strip()
        if not payload:
            continue
        try:
            parsed = json.loads(payload)
            if isinstance(parsed, dict):
                blocks.append(parsed)
        except json.JSONDecodeError:
            continue
    return blocks


def generate_structure(
    prompt: str,
    slide_count: int,
    tone: str,
    style: str,
    presentation_type: str,
    duration_minutes: int,
    context_text: str = "",
    slidesgo_style: str = "",
) -> Dict[str, Any]:
    slidesgo_prompt = f"Шаблон Slidesgo: {slidesgo_style}.\n" if slidesgo_style else ""
    user_message = (
        f"Создай структуру презентации на {slide_count} слайдов.\n"
        f"Тип презентации: {presentation_mode_label(presentation_type)}.\n"
        f"Расчетная длительность: {duration_minutes} мин.\n"
        f"Тон: {tone}.\n"
        f"Визуальный стиль: {style}.\n"
        f"{slidesgo_prompt}\n"
        f"Основной запрос пользователя:\n{prompt}\n\n"
        f"Текст из документа:\n{context_text[:MAX_CONTEXT_CHARS]}"
    )
    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": LLM_MODEL,
            "system_prompt": build_system_prompt(slide_count, tone, style, presentation_type, slidesgo_style),
            "max_new_tokens": 1536,
            "temperature": 0.2,
            "top_p": 0.9,
            "chat_history": [
                {"role": "system", "content": build_system_prompt(slide_count, tone, style, presentation_type, slidesgo_style)},
                {"role": "user", "content": user_message},
            ],
        },
    }

    response = API_SESSION.post(LLM_URL, json=payload, headers=llm_headers(), timeout=REQUEST_TIMEOUT)
    response.raise_for_status()
    data = response.json()

    if isinstance(data, dict) and "slides" in data:
        return normalize_structure(data, slide_count)

    raw_content = ""
    if isinstance(data, list) and data:
        raw_content = data[0].get("message", {}).get("content", "")

    parsed = extract_json_object(raw_content)

    return normalize_structure(parsed, slide_count)


def format_request_error(prefix: str, err: requests.RequestException) -> str:
    err_text = str(err)
    if "IncompleteRead" in err_text or isinstance(getattr(err, "__cause__", None), IncompleteRead):
        return (
            f"{prefix}: соединение оборвалось во время получения ответа от AI сервиса. "
            "Попробуйте повторить запрос."
        )
    if isinstance(err, requests.ConnectTimeout):
        return (
            f"{prefix}: таймаут подключения к внешнему AI сервису. "
            "Проверьте интернет/доступ к ai.rt.ru и попробуйте повторить запрос."
        )
    if isinstance(err, requests.ReadTimeout):
        return (
            f"{prefix}: сервис ответил слишком медленно (таймаут ожидания). "
            "Попробуйте еще раз или уменьшите объем запроса."
        )
    response = getattr(err, "response", None)
    details = ""
    if response is not None:
        try:
            details = response.text.strip()
        except Exception:
            details = ""
    if details:
        return f"{prefix}: {response.status_code} {details}"
    return f"{prefix}: {err}"


def image_download_url(image_id: str, image_type: str = "png") -> str:
    return f"{IMAGE_DOWNLOAD_URL}?id={image_id}&serviceType=yaArt&imageType={image_type}"


def image_preview_url(image_id: str, image_type: str = "png") -> str:
    if not image_id:
        return ""
    return f"/api/image/preview?id={image_id}&image_type={image_type}"


def download_image_with_retries(image_id: str, image_type: str = "png", save_dir: Optional[Path] = None) -> Optional[Path]:
    for attempt in range(1, max(1, IMAGE_PREVIEW_RETRIES) + 1):
        path = download_image(image_id=image_id, image_type=image_type, save_dir=save_dir)
        if path and path.exists():
            return path
        if attempt < IMAGE_PREVIEW_RETRIES:
            sleep(IMAGE_PREVIEW_RETRY_DELAY)
    return None


def generate_image(prompt: str, aspect: str = "16:9") -> Dict[str, str]:
    if not prompt.strip():
        return {"image_id": "", "image_url": ""}

    payload = {
        "uuid": str(uuid.uuid4()),
        "image": {
            "request": prompt,
            "seed": random.randint(1, 2_147_483_647),
            "translate": False,
            "model": "yandex-art",
            "aspect": aspect,
        },
    }

    last_exception: Optional[requests.RequestException] = None
    for attempt in range(1, max(1, IMAGE_GEN_RETRIES) + 1):
        try:
            response = API_SESSION.post(IMAGE_GEN_URL, json=payload, headers=llm_headers(), timeout=IMAGE_REQUEST_TIMEOUT)
            response.raise_for_status()
            data = response.json()
            if not isinstance(data, list) or not data:
                raise RuntimeError("Yandex ART вернул неожиданный формат ответа.")

            message = data[0].get("message", {})
            image_id = str(message.get("id", "")).strip()
            if not image_id:
                raise RuntimeError("Yandex ART не вернул id изображения.")

            local_path = download_image_with_retries(image_id=image_id, image_type="png", save_dir=GENERATED_IMAGE_DIR)
            image_url = f"/static/generated_images/{local_path.name}" if local_path else image_preview_url(image_id)
            return {"image_id": image_id, "image_url": image_url}
        except requests.RequestException as err:
            last_exception = err
            if attempt >= IMAGE_GEN_RETRIES:
                raise
            sleep(IMAGE_GEN_BACKOFF * attempt)

    if last_exception is not None:
        raise last_exception

    raise RuntimeError("Не удалось получить ответ от сервиса генерации изображений.")


def populate_slide_images(structure: Dict[str, Any], aspect: str = "16:9") -> Dict[str, Any]:
    slides = structure.get("slides", [])
    if not isinstance(slides, list):
        return structure

    for slide in slides:
        if not isinstance(slide, dict):
            continue
        image_prompt = str(slide.get("image_prompt", "")).strip()
        image_enabled = bool(slide.get("image_enabled", True)) if image_prompt else False
        image_count = max(1, min(int(slide.get("image_count", 1)), 3))
        slide["image_ids"] = []
        slide["image_urls"] = []
        if not image_prompt or not image_enabled:
            slide["image_id"] = ""
            slide["image_url"] = ""
            continue
        for _ in range(image_count):
            try:
                image_data = generate_image(image_prompt, aspect=aspect)
                slide["image_ids"].append(image_data["image_id"])
                slide["image_urls"].append(image_data["image_url"])
            except (requests.RequestException, RuntimeError):
                break
        if slide["image_ids"]:
            slide["image_id"] = slide["image_ids"][0]
            slide["image_url"] = slide["image_urls"][0]
        else:
            slide["image_id"] = ""
            slide["image_url"] = ""

    return structure


def download_image(image_url: str = "", image_id: str = "", image_type: str = "png", save_dir: Optional[Path] = None) -> Optional[Path]:
    url = image_download_url(image_id, image_type) if image_id else image_url
    if not url:
        return None
    try:
        response = API_SESSION.get(url, headers={"Authorization": f"Bearer {API_TOKEN}"}, timeout=IMAGE_DOWNLOAD_TIMEOUT)
        response.raise_for_status()
        target_dir = save_dir or WORK_DIR
        target_dir.mkdir(parents=True, exist_ok=True)
        path = target_dir / f"{uuid.uuid4().hex}.{image_type}"
        path.write_bytes(response.content)
        return path
    except requests.RequestException:
        return None


def download_images(image_ids: List[str], image_type: str = "png") -> List[Path]:
    image_paths: List[Path] = []
    for image_id in image_ids:
        if not image_id:
            continue
        try:
            response = API_SESSION.get(
                image_download_url(image_id, image_type),
                headers={"Authorization": f"Bearer {API_TOKEN}"},
                timeout=IMAGE_DOWNLOAD_TIMEOUT,
            )
            response.raise_for_status()
            path = WORK_DIR / f"{uuid.uuid4().hex}.{image_type}"
            path.write_bytes(response.content)
            image_paths.append(path)
        except requests.RequestException:
            continue
    return image_paths


def get_style_config(style_key: str, overrides: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    config = STYLE_PRESETS.get(style_key, STYLE_PRESETS["professional"]).copy()
    if style_key in SLIDESGO_STYLES_DESIGN:
        slidesgo_config = SLIDESGO_STYLES_DESIGN[style_key].copy()
        config.update(slidesgo_config)
    if overrides:
        config.update(overrides)
    return config


def apply_style(slide, palette: Dict[str, Any], preserve_background: bool = False) -> None:
    if not preserve_background:
        fill = slide.background.fill
        fill.solid()
        bg_color = palette["bg"]
        if isinstance(bg_color, str):
            bg_color = BACKGROUND_OPTIONS.get(bg_color, RGBColor(255, 255, 255))
        fill.fore_color.rgb = bg_color

    font_family = palette.get("font_family", "Arial")

    if slide.shapes.title and slide.shapes.title.text_frame and slide.shapes.title.text_frame.paragraphs:
        p = slide.shapes.title.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.color.rgb = palette.get("title", RGBColor(0, 0, 0))
            p.runs[0].font.bold = True
            p.runs[0].font.size = Pt(palette.get("title_size", 32))
            p.runs[0].font.name = font_family

    # Apply to body text
    for shape in slide.shapes:
        if shape != slide.shapes.title and hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = palette.get("body", RGBColor(0, 0, 0))
                    run.font.size = Pt(palette.get("body_size", 18))
                    run.font.name = font_family


def fit_image_to_box(image_path: Path, max_width: int, max_height: int) -> Dict[str, int]:

    pixmap = fitz.Pixmap(str(image_path))
    image_width = max(pixmap.width, 1)
    image_height = max(pixmap.height, 1)

    scale = min(max_width / image_width, max_height / image_height)
    width = max(1, int(image_width * scale))
    height = max(1, int(image_height * scale))

    return {"width": width, "height": height}


def build_notes_prompt(structure: Dict[str, Any], prompt: str, presentation_type: str, duration_minutes: int) -> str:
    slides = structure.get("slides", [])
    slide_lines: List[str] = []
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        slide_lines.append(
            f"Слайд {index}\n"
            f"Заголовок: {slide.get('title', '')}\n"
            f"Содержание: {slide.get('content', '')}\n"
        )

    return (
        "Подготовь шпоры для выступающего по презентации.\n"
        f"Тип презентации: {presentation_mode_label(presentation_type)}.\n"
        f"Общая длительность: {duration_minutes} минут.\n"
        f"Исходная тема: {prompt}\n\n"
        "Верни обычный текст, без JSON. Для каждого слайда дай:\n"
        "- номер и заголовок;\n"
        "- что сказать устно;\n"
        "- ключевой акцент;\n"
        "- возможный переход к следующему слайду.\n\n"
        "Структура слайдов:\n"
        + "\n".join(slide_lines)
    )


def generate_speaker_notes(structure: Dict[str, Any], prompt: str, presentation_type: str, duration_minutes: int) -> str:
    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": LLM_MODEL,
            "system_prompt": (
                "Ты — помощник спикера. Пишешь краткие, полезные и естественные шпоры для выступления по слайдам."
            ),
            "max_new_tokens": 1800,
            "temperature": 0.3,
            "top_p": 0.9,
            "chat_history": [
                {
                    "role": "system",
                    "content": "Ты — помощник спикера. Пишешь краткие, полезные и естественные шпоры для выступления по слайдам.",
                },
                {
                    "role": "user",
                    "content": build_notes_prompt(structure, prompt, presentation_type, duration_minutes),
                },
            ],
        },
    }

    response = API_SESSION.post(LLM_URL, json=payload, headers=llm_headers(), timeout=REQUEST_TIMEOUT)
    response.raise_for_status()
    data = response.json()

    if isinstance(data, list) and data:
        return str(data[0].get("message", {}).get("content", "")).strip()
    if isinstance(data, dict):
        return str(data.get("content", "")).strip()
    return ""


def build_notes_json_prompt(structure: Dict[str, Any], prompt: str, presentation_type: str, duration_minutes: int) -> str:
    slides = structure.get("slides", [])
    slide_lines: List[str] = []
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        slide_lines.append(
            f"Слайд {index}\n"
            f"Заголовок: {slide.get('title', '')}\n"
            f"Содержание: {slide.get('content', '')}\n"
        )
    return (
        "Подготовь краткие заметки спикера по каждому слайду.\n"
        f"Тип презентации: {presentation_mode_label(presentation_type)}.\n"
        f"Общая длительность: {duration_minutes} минут.\n"
        f"Исходная тема: {prompt}\n\n"
        "Верни строго JSON формата:\n"
        '{"notes":[{"slide_number":1,"speaker_notes":"..."}]}\n'
        "Заметки должны быть практичными и читаемыми для устного выступления.\n"
        "Структура слайдов:\n"
        + "\n".join(slide_lines)
    )


def generate_speaker_notes_by_slide(structure: Dict[str, Any], prompt: str, presentation_type: str, duration_minutes: int) -> List[str]:
    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": LLM_MODEL,
            "system_prompt": (
                "Ты — помощник спикера. Пишешь короткие и полезные заметки к каждому слайду. "
                "Всегда отвечаешь строго валидным JSON."
            ),
            "max_new_tokens": 2200,
            "temperature": 0.3,
            "top_p": 0.9,
            "chat_history": [
                {
                    "role": "system",
                    "content": (
                        "Ты — помощник спикера. Пишешь короткие и полезные заметки к каждому слайду. "
                        "Всегда отвечаешь строго валидным JSON."
                    ),
                },
                {
                    "role": "user",
                    "content": build_notes_json_prompt(structure, prompt, presentation_type, duration_minutes),
                },
            ],
        },
    }

    response = API_SESSION.post(LLM_URL, json=payload, headers=llm_headers(), timeout=REQUEST_TIMEOUT)
    response.raise_for_status()
    data = response.json()

    content = ""
    if isinstance(data, list) and data:
        content = str(data[0].get("message", {}).get("content", "")).strip()
    elif isinstance(data, dict):
        content = str(data.get("content", "")).strip()

    try:
        parsed = extract_json_object(content)
        raw_notes = parsed.get("notes", [])
        if not isinstance(raw_notes, list):
            raw_notes = []
    except RuntimeError:
        # Fallback: LLM sometimes returns plain text instead of JSON.
        plain_text = generate_speaker_notes(structure, prompt, presentation_type, duration_minutes)
        return parse_plain_notes_by_slide(plain_text, len(structure.get("slides", [])))

    notes_by_slide: Dict[int, str] = {}
    for item in raw_notes:
        if not isinstance(item, dict):
            continue
        try:
            slide_number = int(item.get("slide_number", 0))
        except (TypeError, ValueError):
            slide_number = 0
        note = str(item.get("speaker_notes", "")).strip()
        if slide_number > 0 and note:
            notes_by_slide[slide_number] = note

    slides = structure.get("slides", [])
    result: List[str] = []
    for index in range(len(slides)):
        result.append(notes_by_slide.get(index + 1, ""))
    return result


def parse_plain_notes_by_slide(text: str, slide_count: int) -> List[str]:
    if slide_count <= 0:
        return []
    cleaned = (text or "").strip()
    if not cleaned:
        return [""] * slide_count

    pattern = re.compile(r"(?:^|\n)\s*Слайд\s+(\d+)\s*[:\-]?\s*", flags=re.IGNORECASE)
    matches = list(pattern.finditer(cleaned))
    if not matches:
        first = cleaned[:2000]
        return [first] + [""] * (slide_count - 1)

    notes_by_index: Dict[int, str] = {}
    for i, match in enumerate(matches):
        try:
            slide_number = int(match.group(1))
        except (TypeError, ValueError):
            continue
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(cleaned)
        chunk = cleaned[start:end].strip()
        if slide_number > 0 and chunk:
            notes_by_index[slide_number] = chunk

    return [notes_by_index.get(i + 1, "") for i in range(slide_count)]


def generate_note_for_single_slide(
    title: str,
    content: str,
    prompt: str,
    presentation_type: str,
    duration_minutes: int,
    slide_number: int,
) -> str:
    clean_content = strip_data_blocks(content)
    user_message = (
        "Подготовь заметку спикера для ОДНОГО слайда.\n"
        f"Тип презентации: {presentation_mode_label(presentation_type)}.\n"
        f"Общая длительность презентации: {duration_minutes} минут.\n"
        f"Номер слайда: {slide_number}.\n"
        f"Тема: {prompt}\n"
        f"Заголовок слайда: {title}\n"
        f"Тезисы слайда: {clean_content}\n\n"
        "Верни только текст заметки для выступления (без JSON, без заголовков)."
    )
    payload = {
        "uuid": str(uuid.uuid4()),
        "chat": {
            "model": LLM_MODEL,
            "system_prompt": "Ты — помощник спикера. Пишешь емкую, понятную заметку для устного выступления.",
            "max_new_tokens": 500,
            "temperature": 0.35,
            "top_p": 0.9,
            "chat_history": [
                {
                    "role": "system",
                    "content": "Ты — помощник спикера. Пишешь емкую, понятную заметку для устного выступления.",
                },
                {"role": "user", "content": user_message},
            ],
        },
    }
    response = API_SESSION.post(LLM_URL, json=payload, headers=llm_headers(), timeout=REQUEST_TIMEOUT)
    response.raise_for_status()
    data = response.json()
    if isinstance(data, list) and data:
        return str(data[0].get("message", {}).get("content", "")).strip()
    if isinstance(data, dict):
        return str(data.get("content", "")).strip()
    return ""


def notes_text_from_structure(structure: Dict[str, Any]) -> str:
    slides = structure.get("slides", [])
    lines: List[str] = []
    for index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        title = str(slide.get("title", "")).strip() or f"Слайд {index}"
        notes = str(slide.get("speaker_notes", "")).strip()
        lines.append(f"Слайд {index}: {title}")
        lines.append(notes or "Заметки не заполнены.")
        lines.append("")
    return "\n".join(lines).strip() or "Шпоры не были сгенерированы."


def clear_presentation_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def pick_body_shape(slide):
    body_shape = None
    for placeholder in slide.placeholders:
        if slide.shapes.title is not None and placeholder == slide.shapes.title:
            continue
        if getattr(placeholder, "has_text_frame", False):
            body_shape = placeholder
            break
    if body_shape is None and len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
    if body_shape is None:
        body_shape = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.6), Inches(4.8))
        if getattr(body_shape, "fill", None) is not None:
            body_shape.fill.background()
        if getattr(body_shape, "line", None) is not None and getattr(body_shape.line, "fill", None) is not None:
            body_shape.line.fill.background()
    return body_shape


def _remove_shape(shape):
    try:
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)
    except Exception:
        pass


def replace_template_images(slide, image_paths):
    if not image_paths:
        return image_paths

    picture_targets = []
    for shape in list(slide.shapes):
        ph_type = None
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            ph_type = None
        if ph_type == 18:
            picture_targets.append(shape)

    if not picture_targets:
        picture_targets = [shape for shape in list(slide.shapes) if getattr(shape, "shape_type", None) == 13]

    remaining_images = list(image_paths)
    for target in picture_targets:
        if not remaining_images:
            _remove_shape(target)
            continue
        image_path = remaining_images.pop(0)
        left = int(target.left)
        top = int(target.top)
        width = int(target.width)
        height = int(target.height)
        _remove_shape(target)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    return remaining_images


def replace_template_placeholders(slide, title: str, content: str):
    title_text = str(title or "")
    content_text = str(content or "")
    if slide.shapes.title and getattr(slide.shapes.title, "text_frame", None):
        set_text_frame_text(slide.shapes.title.text_frame, title_text, preserve_style=True)
    for placeholder in slide.placeholders:
        if placeholder == slide.shapes.title:
            continue
        if not getattr(placeholder, "has_text_frame", False):
            continue
        phf = getattr(placeholder, "placeholder_format", None)
        ph_type = getattr(phf, "type", None)
        if ph_type == 4:  # SUBTITLE
            set_text_frame_text(placeholder.text_frame, content_text, preserve_style=True)
        elif ph_type == 2 and content_text:
            set_text_frame_text(placeholder.text_frame, content_text, preserve_style=True)


def choose_slide_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        has_body = False
        for placeholder in layout.placeholders:
            phf = getattr(placeholder, "placeholder_format", None)
            if phf and getattr(phf, "type", None) == 2:  # BODY
                has_body = True
                break
        if has_body:
            return layout
    return prs.slide_layouts[0] if prs.slide_layouts else None


def set_text_frame_text(text_frame, text: str, preserve_style: bool = False) -> None:
    if text is None:
        text = ""
    default_font = None
    default_alignment = None
    if preserve_style and text_frame.paragraphs:
        first_paragraph = text_frame.paragraphs[0]
        default_alignment = first_paragraph.alignment
        if first_paragraph.runs:
            default_font = first_paragraph.runs[0].font
    text_frame.clear()
    lines = str(text).split("\n") if text else [""]
    for index, line in enumerate(lines):
        if index == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.text = line
        if preserve_style and default_font and paragraph.runs:
            run = paragraph.runs[0]
            if default_font.name:
                run.font.name = default_font.name
            if default_font.size:
                run.font.size = default_font.size
            if getattr(default_font.color, "rgb", None):
                run.font.color.rgb = default_font.color.rgb
        if preserve_style and default_alignment is not None:
            paragraph.alignment = default_alignment


def pick_picture_placeholder(slide):
    for placeholder in slide.placeholders:
        phf = getattr(placeholder, "placeholder_format", None)
        if phf and getattr(phf, "type", None) == 18:  # PP_PLACEHOLDER.PICTURE
            return placeholder
    return None


def apply_adaptive_text_size(text_frame, box_width_emu: int, box_height_emu: int, template_mode: bool) -> None:
    text_len = len((text_frame.text or "").strip())
    if text_len == 0:
        return
    area = max(1, box_width_emu * box_height_emu)
    density = text_len / area
    # Heuristic: templates usually have tighter text boxes, so start a bit smaller.
    if density > 2.4e-7:
        size = 12
    elif density > 1.8e-7:
        size = 14
    elif density > 1.2e-7:
        size = 16
    else:
        size = 18 if not template_mode else 16
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def clamp_image_box_to_slide(
    left: int, top: int, width: int, height: int, slide_width: int, slide_height: int
) -> Dict[str, int]:
    safe_left = max(0, left)
    safe_top = max(0, top)
    max_width = max(1, slide_width - safe_left)
    max_height = max(1, slide_height - safe_top)
    safe_width = max(1, min(width, max_width))
    safe_height = max(1, min(height, max_height))
    return {"left": safe_left, "top": safe_top, "width": safe_width, "height": safe_height}


def add_chart_to_slide(slide, data: Dict[str, Any], left, top, width, height):
    chart_type = str(data.get("type", "bar")).strip().lower()
    labels = data.get("labels", [])
    values = data.get("values", [])
    title = str(data.get("title", "")).strip()
    if not isinstance(labels, list) or not isinstance(values, list) or not labels:
        return
    chart_data = CategoryChartData()
    chart_data.categories = [str(label) for label in labels]
    series_values = []
    for value in values:
        try:
            series_values.append(float(value))
        except (TypeError, ValueError):
            series_values.append(0.0)
    chart_data.add_series(title or chart_type.capitalize(), series_values)
    if chart_type in {"line", "statistical", "histogram"}:
        chart_type_enum = XL_CHART_TYPE.LINE if chart_type == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    elif chart_type == "pie":
        chart_type_enum = XL_CHART_TYPE.PIE
    else:
        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data).chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title


def add_table_to_slide(slide, data: Dict[str, Any], left, top, width, height):
    rows = data.get("rows", [])
    if not isinstance(rows, list) or not rows:
        return
    row_count = len(rows)
    col_count = max((len(row) for row in rows if isinstance(row, list)), default=0)
    if row_count < 1 or col_count < 1:
        return
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table
    column_width = int(width // col_count) if col_count else width
    for col in range(col_count):
        table.columns[col].width = column_width
    for row_index, row_data in enumerate(rows):
        if not isinstance(row_data, list):
            row_data = [str(row_data)]
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.text = str(row_data[col_index]) if col_index < len(row_data) else ""


def render_data_blocks(slide, data_blocks: List[Dict[str, Any]], left, top, width, height):
    if not data_blocks:
        return
    block_height = int(height / max(1, len(data_blocks)))
    current_top = top
    for block in data_blocks:
        if not isinstance(block, dict):
            continue
        block_type = str(block.get("type", "")).strip().lower()
        if block_type == "table":
            add_table_to_slide(slide, block, left, current_top, width, block_height)
        else:
            add_chart_to_slide(slide, block, left, current_top, width, block_height)
        current_top += block_height


def create_pptx(
    structure: Dict[str, Any],
    style: str = "professional",
    slidesgo_style_key: str = "",
    overrides: Optional[Dict[str, Any]] = None,
    pptx_template_key: str = "",
) -> Path:
    template_key = (pptx_template_key or "").strip().lower()
    template_path = Path(app.root_path) / "static" / "templates" / f"style_{template_key}.pptx"
    
    if template_key and template_path.exists():
        # Load template and work directly with its slide structure
        prs = Presentation(str(template_path))
        using_template = True
        # Trim to number of slides needed, or add more if necessary
        num_template_slides = len(prs.slides)
        num_needed_slides = len(structure.get("slides", []))
        
        # Remove excess slides if template has more
        if num_template_slides > num_needed_slides:
            slides_to_remove = num_template_slides - num_needed_slides
            for _ in range(slides_to_remove):
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
    else:
        prs = Presentation()
        using_template = False
    
    image_files: List[Path] = []
    palette = get_style_config(slidesgo_style_key or style, overrides)

    slide_layout = None
    if not using_template:
        slide_layout = choose_slide_layout(prs) or (prs.slide_layouts[0] if prs.slide_layouts else None)
        if slide_layout is None:
            raise RuntimeError("No slide layouts available for PPTX generation.")

    for slide_index, slide_data in enumerate(structure["slides"]):
        if using_template and slide_index < len(prs.slides):
            # Use existing template slide
            slide = prs.slides[slide_index]
        else:
            # Create new slide
            slide = prs.slides.add_slide(slide_layout)
            if not using_template:
                apply_style(slide, palette)

        if using_template:
            replace_template_placeholders(slide, slide_data.get("title", ""), str(slide_data.get("content", "")).strip())
        elif slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        body_shape = pick_body_shape(slide)
        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        margin_left = int(Inches(0.7))
        margin_right = int(Inches(0.7))
        margin_top = int(Inches(1.8))
        margin_bottom = int(Inches(0.7))
        image_spacing = int(Inches(0.25))
        image_outer_max_width = int(Inches(3.6))
        text_box_height = int(Inches(4.8))

        if not using_template:
            body_shape.left = margin_left
            body_shape.top = margin_top
            body_shape.height = text_box_height
            body_shape.width = slide_width - margin_left - margin_right

        image_paths = download_images(slide_data.get("image_ids", []) if isinstance(slide_data.get("image_ids", []), list) else [])
        if not image_paths and slide_data.get("image_id"):
            fallback_path = download_image(slide_data.get("image_url", ""), slide_data.get("image_id", ""))
            if fallback_path:
                image_paths = [fallback_path]

        if using_template:
            image_paths = replace_template_images(slide, image_paths)

        picture_placeholder = pick_picture_placeholder(slide) if image_paths and using_template else None
        image_box_left = None
        image_box_top = None
        image_box_width = None
        image_box_height = None

        if image_paths:
            image_files.extend(image_paths)
            if picture_placeholder is not None:
                image_box_left = int(picture_placeholder.left)
                image_box_top = int(picture_placeholder.top)
                image_box_width = int(picture_placeholder.width)
                image_box_height = int(picture_placeholder.height)
            else:
                available_width = slide_width - margin_left - margin_right - image_spacing
                desired_image_width = min(image_outer_max_width, int(Inches(4.0)))
                reserve_text_width = int(Inches(3.0))
                if available_width >= desired_image_width + reserve_text_width:
                    body_shape.width = max(reserve_text_width, available_width - desired_image_width)
                    image_box_left = body_shape.left + body_shape.width + image_spacing
                    image_box_top = body_shape.top
                    image_box_width = min(desired_image_width, slide_width - margin_right - image_box_left)
                    image_box_height = max(1, min(text_box_height, slide_height - image_box_top - margin_bottom))
                else:
                    # если справа не помещается, спускаем картинку ниже текста
                    body_shape.width = slide_width - margin_left - margin_right
                    image_box_left = margin_left
                    image_box_top = body_shape.top + body_shape.height + image_spacing
                    image_box_width = slide_width - margin_left - margin_right
                    image_box_height = max(1, slide_height - image_box_top - margin_bottom)
                    if image_box_height < int(Inches(2.0)):
                        image_box_height = max(1, slide_height - margin_top - margin_bottom)
                        image_box_top = margin_top
                        image_box_left = slide_width - margin_right - min(desired_image_width, available_width)
                        image_box_width = min(desired_image_width, available_width)
                        if image_box_width < int(Inches(2.0)):
                            image_box_left = margin_left
                            image_box_width = slide_width - margin_left - margin_right
                            image_box_top = body_shape.top + body_shape.height + image_spacing
                            image_box_height = max(1, slide_height - image_box_top - margin_bottom)
                if using_template and image_box_left is None:
                    body_shape.width = max(int(Inches(3.5)), min(body_shape.width, slide_width - body_shape.left - margin_right - image_spacing - int(Inches(3.0))))
                    image_box_left = body_shape.left + body_shape.width + image_spacing
                    image_box_top = body_shape.top if hasattr(body_shape, 'top') else margin_top
                    image_box_width = max(1, min(image_outer_max_width, slide_width - margin_right - image_box_left))
                    image_box_height = max(1, min(body_shape.height if hasattr(body_shape, 'height') else text_box_height, slide_height - image_box_top - margin_bottom))

        content_text = str(slide_data.get("content", "")).strip()
        data_blocks = slide_data.get("data_blocks", []) if isinstance(slide_data.get("data_blocks", []), list) else []
        if data_blocks and not using_template:
            body_shape.height = int(Inches(3.0))

        set_text_frame_text(body_shape.text_frame, content_text, preserve_style=bool(using_template))
        tf = body_shape.text_frame
        tf.word_wrap = True
        apply_adaptive_text_size(tf, int(body_shape.width), int(body_shape.height), bool(using_template))
        if not using_template:
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = palette.get("body", RGBColor(55, 65, 81))
                    run.font.size = Pt(palette.get("body_size", 18))
                    run.font.name = palette.get("font_family", "Calibri")

        if data_blocks:
            render_data_blocks(
                slide,
                data_blocks,
                int(body_shape.left),
                int(body_shape.top + body_shape.height + Inches(0.2)),
                int(body_shape.width),
                int(Inches(3.5)),
            )

        if image_paths:
            safe_box = clamp_image_box_to_slide(
                int(image_box_left),
                int(image_box_top),
                int(image_box_width),
                int(image_box_height),
                int(prs.slide_width),
                int(prs.slide_height),
            )
            if len(image_paths) == 1:
                image_size = fit_image_to_box(image_paths[0], safe_box["width"], safe_box["height"])
                image_left = safe_box["left"] + int((safe_box["width"] - image_size["width"]) / 2)
                image_top = safe_box["top"] + int((safe_box["height"] - image_size["height"]) / 2)
                image_left = max(0, min(image_left, int(prs.slide_width) - image_size["width"]))
                image_top = max(0, min(image_top, int(prs.slide_height) - image_size["height"]))
                slide.shapes.add_picture(
                    str(image_paths[0]),
                    image_left,
                    image_top,
                    width=image_size["width"],
                    height=image_size["height"],
                )
            else:
                columns = min(3, len(image_paths))
                gap = int(Inches(0.15))
                image_width = max(1, int((safe_box["width"] - gap * (columns - 1)) / columns))
                image_height = safe_box["height"]
                for i, image_path in enumerate(image_paths[:columns]):
                    image_size = fit_image_to_box(image_path, image_width, image_height)
                    image_left = safe_box["left"] + i * (image_width + gap)
                    image_top = safe_box["top"] + int((safe_box["height"] - image_size["height"]) / 2)
                    slide.shapes.add_picture(
                        str(image_path),
                        image_left,
                        image_top,
                        width=image_size["width"],
                        height=image_size["height"],
                    )

    output_path = WORK_DIR / f"presentation_{uuid.uuid4().hex}.pptx"
    prs.save(str(output_path))

    for image_file in image_files:
        try:
            image_file.unlink(missing_ok=True)
        except OSError:
            pass

    return output_path


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/slidesgo_styles")
def api_slidesgo_styles():
    try:
        return jsonify(get_slidesgo_styles())
    except Exception as err:
        return jsonify({"categories": [{"key": "all", "title": "Все стили"}], "styles": {"all": []}, "error": str(err)}), 502


@app.route("/api/structure", methods=["POST"])
def api_structure():
    prompt = request.form.get("prompt", "").strip()
    tone = request.form.get("tone", "professional")
    style = request.form.get("style", "professional")
    presentation_type = request.form.get("presentation_type", "public").strip() or "public"
    count_mode = request.form.get("count_mode", "minutes").strip() or "minutes"
    duration_raw = request.form.get("duration_minutes", "5")
    slide_count_raw = request.form.get("slide_count", "")
    generate_images = request.form.get("generate_images", "").lower() in {"1", "true", "on", "yes"}
    slidesgo_style_key = request.form.get("slidesgo_style", "").strip()
    slidesgo_style_data = get_slidesgo_style_by_key(slidesgo_style_key)
    slidesgo_style = ""
    if slidesgo_style_data:
        slidesgo_style = slidesgo_style_data.get("title", "")
        description = slidesgo_style_data.get("description", "")
        if description:
            slidesgo_style += f" — {description}"

    try:
        duration_minutes = max(1, min(int(duration_raw), 120))
    except ValueError:
        duration_minutes = 5

    if count_mode == "slides":
        try:
            slide_count = max(3, min(int(slide_count_raw), 20))
        except ValueError:
            slide_count = calculate_slide_count(duration_minutes, presentation_type)
        minutes_per_slide = PUBLIC_MINUTES_PER_SLIDE if presentation_type == "public" else SELF_STUDY_MINUTES_PER_SLIDE
        duration_minutes = max(1, min(round(slide_count * minutes_per_slide), 120))
    else:
        slide_count = calculate_slide_count(duration_minutes, presentation_type)

    if not prompt and "file" not in request.files:
        return jsonify({"error": "Введите промпт или загрузите документ."}), 400

    context_text = ""
    uploaded_file = request.files.get("file")
    if uploaded_file and uploaded_file.filename:
        if not allowed_file(uploaded_file.filename):
            return jsonify({"error": "Поддерживаются только PDF, DOCX и TXT."}), 400
        context_text = extract_text(uploaded_file)

    try:
        structure = generate_structure(
            prompt,
            slide_count,
            tone,
            style,
            presentation_type,
            duration_minutes,
            context_text,
            slidesgo_style=slidesgo_style,
        )
        if generate_images:
            structure = populate_slide_images(structure)
        return jsonify({"structure": structure, "slide_count": slide_count})
    except requests.RequestException as err:
        error_prefix = "Ошибка генерации изображения" if generate_images else "Ошибка LLM API"
        return jsonify({"error": format_request_error(error_prefix, err)}), 502
    except RuntimeError as err:
        return jsonify({"error": str(err)}), 500


@app.route("/api/image/regenerate", methods=["POST"])
def api_regenerate_image():
    payload = request.get_json(silent=True) or {}
    image_prompt = str(payload.get("image_prompt", "")).strip()
    if not image_prompt:
        return jsonify({"error": "Пустой image_prompt."}), 400

    try:
        image_data = generate_image(image_prompt)
        return jsonify(image_data)
    except requests.RequestException as err:
        return jsonify({"error": format_request_error("Ошибка генерации изображения", err)}), 502
    except RuntimeError as err:
        return jsonify({"error": str(err)}), 500


@app.route("/api/image/preview")
def api_image_preview():
    image_id = request.args.get("id", "").strip()
    image_type = request.args.get("image_type", "png").strip() or "png"
    if not image_id:
        return jsonify({"error": "Пустой id изображения."}), 400

    tried_errors: List[str] = []
    print(f"[DEBUG] api_image_preview: id={image_id}, type={image_type}, TOKEN={API_TOKEN[:20]}...")
    
    for attempt in range(max(1, IMAGE_PREVIEW_RETRIES)):
        for candidate_type in (image_type, "jpg", "png"):
            try:
                url = image_download_url(image_id, candidate_type)
                print(f"[DEBUG] Attempt {attempt+1}: Fetching {url} with timeout={IMAGE_DOWNLOAD_TIMEOUT}s")
                response = API_SESSION.get(
                    url,
                    headers={"Authorization": f"Bearer {API_TOKEN}"},
                    timeout=IMAGE_DOWNLOAD_TIMEOUT,
                )
                print(f"[DEBUG] Response status: {response.status_code}, content_type: {response.headers.get('content-type')}")
                if response.status_code != 200:
                    print(f"[DEBUG] Non-200 response: {response.text[:200]}")
                response.raise_for_status()
                print(f"[DEBUG] Success! Returning image, size={len(response.content)} bytes")
                return send_file(
                    io.BytesIO(response.content),
                    mimetype=f"image/{candidate_type}",
                    download_name=f"image.{candidate_type}",
                )
            except requests.Timeout as err:
                error_msg = f"Timeout ({IMAGE_DOWNLOAD_TIMEOUT}s) loading image {image_id}"
                print(f"[DEBUG] TIMEOUT: {error_msg}")
                tried_errors.append(error_msg)
                continue
            except requests.HTTPError as err:
                error_msg = f"HTTP {response.status_code}: {response.text[:100]}"
                print(f"[DEBUG] HTTP Error: {error_msg}")
                tried_errors.append(error_msg)
                continue
            except requests.RequestException as err:
                error_msg = str(err)
                print(f"[DEBUG] Request Error: {error_msg}")
                tried_errors.append(error_msg)
                continue
            except Exception as err:
                error_msg = f"Unexpected error: {str(err)}"
                print(f"[DEBUG] Unexpected Error: {error_msg}")
                tried_errors.append(error_msg)
                continue
        if attempt < IMAGE_PREVIEW_RETRIES - 1:
            print(f"[DEBUG] Sleeping {IMAGE_PREVIEW_RETRY_DELAY}s before next attempt...")
            sleep(IMAGE_PREVIEW_RETRY_DELAY)
    detail = tried_errors[-1] if tried_errors else "unknown error"
    print(f"[DEBUG] All attempts failed: {detail}")
    return jsonify({"error": f"Ошибка загрузки изображения: {detail}"}), 502


@app.route("/api/pptx", methods=["POST"])
def api_pptx():
    payload = request.get_json(silent=True) or {}
    structure = payload.get("structure", {})
    style = str(payload.get("style", "professional"))
    slidesgo_style_key = str(payload.get("slidesgo_style", "")).strip()
    pptx_template_key = str(payload.get("pptx_template", "")).strip()
    background = str(payload.get("background", "white"))
    font_family = str(payload.get("font_family", "Arial"))
    layout = str(payload.get("layout", "centered"))
    image_position = str(payload.get("image_position", "center"))

    if not isinstance(structure, dict) or "slides" not in structure:
        return jsonify({"error": "Неверный JSON структуры презентации."}), 400

    overrides = {
        "bg": background,
        "font_family": font_family,
        "layout": layout,
        "image_position": image_position,
    }

    try:
        safe_structure = normalize_structure(structure, len(structure.get("slides", [])) or 6)
        if any(
            isinstance(slide, dict)
            and bool(slide.get("image_enabled", False))
            and str(slide.get("image_prompt", "")).strip()
            and not slide.get("image_id")
            for slide in safe_structure["slides"]
        ):
            safe_structure = populate_slide_images(safe_structure)
        pptx_path = create_pptx(
            safe_structure,
            style=style,
            slidesgo_style_key=slidesgo_style_key,
            overrides=overrides,
            pptx_template_key=pptx_template_key,
        )
        return send_file(
            str(pptx_path),
            as_attachment=True,
            download_name="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as err:
        return jsonify({"error": f"Ошибка сборки PPTX: {err}"}), 500


@app.route("/api/notes", methods=["POST"])
def api_notes():
    payload = request.get_json(silent=True) or {}
    structure = payload.get("structure", {})
    prompt = str(payload.get("prompt", "")).strip()
    presentation_type = str(payload.get("presentation_type", "public")).strip() or "public"

    try:
        duration_minutes = max(1, min(int(payload.get("duration_minutes", 5)), 120))
    except (TypeError, ValueError):
        duration_minutes = 5

    if not isinstance(structure, dict) or "slides" not in structure:
        return jsonify({"error": "Неверный JSON структуры презентации."}), 400

    try:
        safe_structure = normalize_structure(structure, len(structure.get("slides", [])) or 6)
        source_slides = structure.get("slides", []) if isinstance(structure, dict) else []
        for index, slide in enumerate(safe_structure.get("slides", [])):
            if index >= len(source_slides) or not isinstance(source_slides[index], dict):
                continue
            slide["speaker_notes"] = str(source_slides[index].get("speaker_notes", "")).strip()

        has_custom_notes = any(str(slide.get("speaker_notes", "")).strip() for slide in safe_structure.get("slides", []))
        notes_text = notes_text_from_structure(safe_structure) if has_custom_notes else generate_speaker_notes(
            safe_structure, prompt, presentation_type, duration_minutes
        )
        notes_path = WORK_DIR / f"speaker_notes_{uuid.uuid4().hex}.txt"
        notes_path.write_text(notes_text or "Шпоры не были сгенерированы.", encoding="utf-8")
        return send_file(str(notes_path), as_attachment=True, download_name="speaker_notes.txt", mimetype="text/plain; charset=utf-8")
    except requests.RequestException as err:
        return jsonify({"error": format_request_error("Ошибка генерации шпор", err)}), 502
    except RuntimeError as err:
        return jsonify({"error": str(err)}), 500


@app.route("/api/notes/generate", methods=["POST"])
def api_notes_generate():
    payload = request.get_json(silent=True) or {}
    structure = payload.get("structure", {})
    prompt = str(payload.get("prompt", "")).strip()
    presentation_type = str(payload.get("presentation_type", "public")).strip() or "public"

    try:
        duration_minutes = max(1, min(int(payload.get("duration_minutes", 5)), 120))
    except (TypeError, ValueError):
        duration_minutes = 5

    if not isinstance(structure, dict) or "slides" not in structure:
        return jsonify({"error": "Неверный JSON структуры презентации."}), 400

    try:
        safe_structure = normalize_structure(structure, len(structure.get("slides", [])) or 6)
        notes = generate_speaker_notes_by_slide(safe_structure, prompt, presentation_type, duration_minutes)
        return jsonify({"notes": notes})
    except requests.RequestException as err:
        return jsonify({"error": format_request_error("Ошибка генерации шпор", err)}), 502
    except RuntimeError as err:
        return jsonify({"error": str(err)}), 500


@app.route("/api/notes/regenerate_slide", methods=["POST"])
def api_notes_regenerate_slide():
    payload = request.get_json(silent=True) or {}
    slide = payload.get("slide", {})
    prompt = str(payload.get("prompt", "")).strip()
    presentation_type = str(payload.get("presentation_type", "public")).strip() or "public"
    try:
        duration_minutes = max(1, min(int(payload.get("duration_minutes", 5)), 120))
    except (TypeError, ValueError):
        duration_minutes = 5
    try:
        slide_number = max(1, int(payload.get("slide_number", 1)))
    except (TypeError, ValueError):
        slide_number = 1

    if not isinstance(slide, dict):
        return jsonify({"error": "Неверный формат слайда."}), 400

    title = str(slide.get("title", "")).strip()
    content = str(slide.get("content", "")).strip()
    if not title and not content:
        return jsonify({"error": "Заполните заголовок или тезисы слайда."}), 400

    try:
        note = generate_note_for_single_slide(title, content, prompt, presentation_type, duration_minutes, slide_number)
        return jsonify({"speaker_note": note})
    except requests.RequestException as err:
        return jsonify({"error": format_request_error("Ошибка генерации шпоры по слайду", err)}), 502
    except RuntimeError as err:
        return jsonify({"error": str(err)}), 500


@app.route("/generate-pptx", methods=["POST"])
def generate_pptx_from_template():
    payload = request.get_json(silent=True) or request.form
    title = str(payload.get("title", "")).strip() or "Заголовок презентации"
    subtitle = str(payload.get("subtitle", "")).strip() or "Подзаголовок"
    style = str(payload.get("style", "blobs")).strip().lower() or "blobs"

    template_map = {
        "blobs": "style_blobs.pptx",
        "geometric": "style_geometric.pptx",
    }
    template_name = template_map.get(style)
    if not template_name:
        return jsonify({"error": "Неизвестный стиль шаблона."}), 400

    template_path = Path(app.root_path) / "static" / "templates" / template_name
    if not template_path.exists():
        return jsonify({"error": f"Шаблон не найден: {template_name}"}), 404

    presentation = Presentation(str(template_path))
    if not presentation.slides:
        return jsonify({"error": "Шаблон презентации не содержит слайдов."}), 500

    slide = presentation.slides[0]
    replace_template_placeholders(slide, title, subtitle)

    if style == "geometric":
        triangle_specs = [
            (Inches(-0.2), Inches(-0.2), Inches(2.8), Inches(2.2), RGBColor(33, 99, 235)),   # синий
            (Inches(11.0), Inches(-0.2), Inches(2.5), Inches(2.0), RGBColor(59, 191, 255)),  # голубой
            (Inches(11.1), Inches(6.0), Inches(2.6), Inches(2.1), RGBColor(250, 204, 21)),   # желтый
        ]
        for left, top, width, height, color in triangle_specs:
            triangle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
            triangle.fill.solid()
            triangle.fill.fore_color.rgb = color
            triangle.line.fill.background()

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return send_file(
        output,
        as_attachment=True,
        download_name="presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


if __name__ == "__main__":
    port = int(os.getenv("PORT", "5000"))
    app.run(debug=True, host="0.0.0.0", port=port, threaded=True)
